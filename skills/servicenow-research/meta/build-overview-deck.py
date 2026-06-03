"""Generate the Accenture-branded overview deck for the servicenow-research skill.

14 slides: cover + agenda + 12 content. Patterns: I, D, L, V, P, G, X, Q, Y, A, E, AB, J.
Theme: accenture (white cover, BLACK text). Language: Japanese.
"""
from __future__ import annotations
import sys, os

_SKILL = os.path.join(os.path.expanduser("~"), ".claude", "skills", "acnpptx", "scripts")
sys.path.insert(0, _SKILL)

import helpers as _h
_h.set_lang("ja")
_h.load_theme("accenture")
# Workaround: load_theme("accenture") sets DARK_PURPLE=#7300B9 which is NOT in
# accenture.json tokens. brand_check rejects it. Realign to accenture's
# primary_medium (#7500C0).
from pptx.dml.color import RGBColor as _RGB
_h.DARK_PURPLE = _RGB(0x75, 0x00, 0xC0)
from helpers import *  # noqa: E402,F401,F403

from native_shapes import (  # noqa: E402
    add_chevron_flow, add_arrow_right, add_connector_arrow, add_divider_line,
    add_highlight_bar,
)
from pattern_v import add_numbered_card_grid  # noqa: E402
from pattern_x import add_step_chart  # noqa: E402

from pptx import Presentation  # noqa: E402
from pptx.util import Inches, Pt, Emu  # noqa: E402
from pptx.dml.color import RGBColor  # noqa: E402
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR  # noqa: E402
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR_TYPE  # noqa: E402
from pptx.oxml.ns import qn as _qn  # noqa: E402

# Theme-specific layout indices (accenture)
LAYOUT_COVER = 0
LAYOUT_CONTENT = 2

prs = Presentation(TEMPLATE_PATH)
# Strip any pre-existing slides from the template
while len(prs.slides) > 0:
    sldId = prs.slides._sldIdLst[0]
    prs.part.drop_rel(sldId.get(_qn("r:id")))
    del prs.slides._sldIdLst[0]


# ── Helpers ───────────────────────────────────────────────────────────────

def add_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_CONTENT])
    clear_placeholders(slide)
    return slide


def add_cover_slide_layout():
    return prs.slides.add_slide(prs.slide_layouts[LAYOUT_COVER])


def _set_anchor_ctr(tb):
    bodyPr = tb.text_frame._txBody.find(_qn('a:bodyPr'))
    if bodyPr is not None:
        bodyPr.set('anchor', 'ctr')


def _add_textbox(slide, x, y, w, h, *, anchor_ctr=False, word_wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = word_wrap
    if anchor_ctr:
        _set_anchor_ctr(tb)
    return tb, tf


def _set_para(p, text, *, size=14, bold=False, color=None, font_name=None, align=None,
              space_after_pt=8):
    p.text = text
    if align is not None:
        p.alignment = align
    if space_after_pt is not None:
        p.space_after = Pt(space_after_pt)
    p.font.size = Pt(size)
    p.font.bold = bool(bold)
    if color is not None:
        p.font.color.rgb = color
    p.font.name = font_name or FONT


def _multi_line(tf, lines, *, size=14, bold=False, color=None, font_name=None,
                align=None, space_after_pt=8):
    color = color if color is not None else TEXT_BODY
    for j, ln in enumerate(lines):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        _set_para(p, ln, size=size, bold=bold, color=color, font_name=font_name,
                  align=align, space_after_pt=space_after_pt)


def _add_rect(slide, x, y, w, h, *, fill=None, line_fill_bg=True, line_color=None,
              line_w_pt=None):
    shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is not None:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line_color is not None:
        shp.line.color.rgb = line_color
        if line_w_pt is not None:
            shp.line.width = Pt(line_w_pt)
    elif line_fill_bg:
        shp.line.fill.background()
    return shp


# ───────────────────────────────────────────────────────────────────────────
# Slide 1 — Cover
# ───────────────────────────────────────────────────────────────────────────
def slide_01_cover():
    slide = add_cover_slide_layout()
    # Per accenture theme: idx_0 title, idx_1 subtitle, idx_2 date, idx_12 body
    # All BLACK text on white.
    title_set = False
    for ph in list(slide.placeholders):
        idx = ph.placeholder_format.idx
        ph.text_frame.clear()
        p = ph.text_frame.paragraphs[0]
        if idx == 0:
            p.text = "ServiceNow Research スキル"
            p.font.size = Pt(40); p.font.bold = True
            p.font.color.rgb = BLACK; p.font.name = FONT
            p2 = ph.text_frame.add_paragraph()
            p2.text = "Claude Code 用 公式情報調査ツール"
            p2.font.size = Pt(28); p2.font.bold = False
            p2.font.color.rgb = CORE_PURPLE; p2.font.name = FONT
            title_set = True
        elif idx == 1:
            p.text = "ServiceNowDocs を主軸にした GitHub-First リサーチ自動化"
            p.font.size = Pt(20); p.font.bold = False
            p.font.color.rgb = TEXT_BODY; p.font.name = FONT
        elif idx == 2:
            p.text = "v1.0.0  /  2026年5月"
            p.font.size = Pt(14)
            p.font.color.rgb = MID_GRAY; p.font.name = FONT
        elif idx == 12:
            p.text = "受領者向けスキル概要 — 5分で導入できる構成"
            p.font.size = Pt(14)
            p.font.color.rgb = TEXT_BODY; p.font.name = FONT
        else:
            # Hide unused placeholder hint text
            p.text = " "
            p.font.size = Pt(8); p.font.color.rgb = WHITE; p.font.name = FONT
    if not title_set:
        # Fallback: layout had no title placeholder
        tb, tf = _add_textbox(slide, ML, 2.6, CW, 1.2)
        _set_para(tf.paragraphs[0], "ServiceNow Research スキル",
                  size=40, bold=True, color=BLACK, align=PP_ALIGN.LEFT)


# ───────────────────────────────────────────────────────────────────────────
# Slide 2 — Agenda (Pattern I)
# ───────────────────────────────────────────────────────────────────────────
def slide_02_agenda():
    slide = add_slide()
    add_breadcrumb(slide, "ServiceNow Research > Agenda")
    add_title(slide, "目次")
    # Agenda has no message line per skill convention

    items = [
        ("01", "概要",       "解決する課題 / できること・できないこと"),
        ("02", "アーキテクチャ", "5並列リサーチャ構成 / 処理フロー"),
        ("03", "導入と使い方",   "前提環境 / 3ステップ導入 / 質問例"),
        ("04", "制約と運用",     "対応リリース / 制約事項 / トラブルシュート"),
    ]
    n = len(items)
    row_h = 0.95
    gap = 0.25
    total_h = n * row_h + (n - 1) * gap
    start_y = CY + (AH - total_h) / 2
    badge_w = 1.20
    text_x = ML + badge_w + 0.30
    text_w = CW - badge_w - 0.30

    for i, (num, title_t, desc) in enumerate(items):
        y = start_y + i * (row_h + gap)
        # Badge — purple square
        b = _add_rect(slide, ML, y, badge_w, row_h, fill=CORE_PURPLE)
        tf = b.text_frame
        tf.word_wrap = False  # prevent "0\n1" split
        bodyPr = tf._txBody.find(_qn('a:bodyPr'))
        if bodyPr is not None:
            bodyPr.set('anchor', 'ctr')
            for k in ('lIns', 'rIns', 'tIns', 'bIns'):
                bodyPr.set(k, '0')
        _set_para(tf.paragraphs[0], num, size=44, bold=True, color=WHITE,
                  align=PP_ALIGN.CENTER, space_after_pt=0)
        # Title + description
        tb, tf2 = _add_textbox(slide, text_x, y, text_w, row_h, anchor_ctr=True)
        p1 = tf2.paragraphs[0]
        _set_para(p1, title_t, size=22, bold=True, color=BLACK, space_after_pt=4)
        p2 = tf2.add_paragraph()
        _set_para(p2, desc, size=14, color=TEXT_BODY)

    set_footer(slide)


# ───────────────────────────────────────────────────────────────────────────
# Slide 3 — Solving the problem (Pattern D - Key Message)
# ───────────────────────────────────────────────────────────────────────────
def slide_03_problem():
    slide = add_slide()
    add_breadcrumb(slide, "ServiceNow Research > 概要")
    add_title(slide, "解決する課題")
    add_message_line(slide, "「全部 GitHub に聞ける」体験を Claude Code から提供する")

    # Big centered message
    msg_h = 2.20
    msg_y = CY + 0.40
    panel = _add_rect(slide, ML, msg_y, CW, msg_h, fill=DARKEST_PURPLE)
    tf = panel.text_frame
    tf.word_wrap = True
    _set_anchor_ctr(panel)
    _multi_line(tf, [
        "ServiceNow 公式情報は「サイトをまたいで散在」している",
        "本スキルは GitHub の ServiceNowDocs を主軸に",
        "5並列ソースを自動横断調査し、出典付きで日本語回答する",
    ], size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER, space_after_pt=10)

    # Three pain points below
    pain_y = msg_y + msg_h + 0.30
    pain_h = 1.40
    col_w = (CW - 2 * 0.30) / 3
    pains = [
        ("⚠️ 散在", "docs / developer / community / Now Support に\n情報が分散、横断調査が手作業"),
        ("⚠️ 非効率", "JS SPA / 認証壁で WebFetch 不可、\n手動コピペが日常化"),
        ("⚠️ 誤情報リスク", "Qiita / Stack Overflow など非公式情報の\n混入で仕様の取り違いが発生"),
    ]
    for i, (head, body) in enumerate(pains):
        x = ML + i * (col_w + 0.30)
        _add_rect(slide, x, pain_y, col_w, pain_h, fill=OFF_WHITE)
        tb, tf = _add_textbox(slide, x + 0.20, pain_y + 0.10, col_w - 0.40, pain_h - 0.20,
                              anchor_ctr=True)
        p1 = tf.paragraphs[0]
        _set_para(p1, head, size=16, bold=True, color=CORE_PURPLE, space_after_pt=8)
        for ln in body.split("\n"):
            p = tf.add_paragraph()
            _set_para(p, ln, size=12, color=TEXT_BODY, space_after_pt=4)

    set_footer(slide)


# ───────────────────────────────────────────────────────────────────────────
# Slide 4 — Capabilities (Pattern L - Do/Don't)
# ───────────────────────────────────────────────────────────────────────────
def slide_04_capabilities():
    slide = add_slide()
    add_breadcrumb(slide, "ServiceNow Research > 概要")
    add_title(slide, "できること / できないこと")
    add_message_line(slide, "公式情報の自動調査に特化、操作系・非公式情報には踏み込まない")

    col_w = (CW - 0.40) / 2
    panel_h = AH - 0.40
    panel_y = CY + 0.20
    header_h = 0.55

    # Left: ✅ できる
    _add_rect(slide, ML, panel_y, col_w, header_h, fill=CORE_PURPLE)
    tb, tf = _add_textbox(slide, ML + 0.20, panel_y, col_w - 0.40, header_h, anchor_ctr=True)
    _set_para(tf.paragraphs[0], "✅ できること", size=18, bold=True, color=WHITE,
              space_after_pt=0)
    _add_rect(slide, ML, panel_y + header_h, col_w, panel_h - header_h, fill=OFF_WHITE)
    tb, tf = _add_textbox(slide, ML + 0.25, panel_y + header_h + 0.15,
                          col_w - 0.50, panel_h - header_h - 0.30, anchor_ctr=True)
    _multi_line(tf, [
        "● ServiceNow 公式仕様・API・設定手順の調査",
        "● リリース別 (Australia/Zurich/Yokohama/Xanadu) 差分の取得",
        "● 出典 URL の自動生存確認 (NotFound 排除)",
        "● 専門用語に初心者向け補足を自動付与",
        "● 該当無しは正直に「わかりません」と回答 (捏造防止)",
    ], size=14, color=TEXT_BODY, space_after_pt=10)

    # Right: ❌ できない
    rx = ML + col_w + 0.40
    _add_rect(slide, rx, panel_y, col_w, header_h, fill=MID_GRAY)
    tb, tf = _add_textbox(slide, rx + 0.20, panel_y, col_w - 0.40, header_h, anchor_ctr=True)
    _set_para(tf.paragraphs[0], "❌ できないこと", size=18, bold=True, color=WHITE,
              space_after_pt=0)
    _add_rect(slide, rx, panel_y + header_h, col_w, panel_h - header_h, fill=OFF_WHITE)
    tb, tf = _add_textbox(slide, rx + 0.25, panel_y + header_h + 0.15,
                          col_w - 0.50, panel_h - header_h - 0.30, anchor_ctr=True)
    _multi_line(tf, [
        "● ServiceNow インスタンスの実機操作",
        "● 認証必須の Now Support KB 本文取得 (auth wall)",
        "● Vancouver より古いリリース固有情報 (リポジトリから削除済)",
    ], size=14, color=TEXT_BODY, space_after_pt=12)

    set_footer(slide)


# ───────────────────────────────────────────────────────────────────────────
# Slide 5 — Architecture (Pattern V - Numbered Card Grid)
# ───────────────────────────────────────────────────────────────────────────
def slide_05_architecture():
    slide = add_slide()
    add_breadcrumb(slide, "ServiceNow Research > アーキテクチャ")
    add_title(slide, "5並列リサーチャ構成")
    add_message_line(slide, "99%の質問にGitHubが答え、残り1%を補助ソースが拾う")

    cards = [
        {"number": 1, "title": "GitHub (主)\n~99% ヒット",
         "body": "ServiceNow/ServiceNowDocs リポジトリ\ngh search code + raw URL WebFetch\n4ブランチ (australia/zurich/yokohama/xanadu) 対応"},
        {"number": 2, "title": "Developer\n<10% (SPA壁)",
         "body": "developer.servicenow.com\nWebSearch + WebFetch ハイブリッド\nblog.do / to.do / _escaped_fragment_ のみ採用"},
        {"number": 3, "title": "Docs\n<5% (SPA壁)",
         "body": "docs.servicenow.com\nWebSearch + WebFetch\nGoogle キャッシュ越しに偶然取れる程度"},
        {"number": 4, "title": "Community\n~30% (SSR Q&A)",
         "body": "community.servicenow.com\nQ&A (?id=community_question) は SSR\n実務トラブル解決に有効"},
        {"number": 5, "title": "Now Support\n~0% (認証壁)",
         "body": "support.servicenow.com\nKB 本文は SSO 必須\n設計上 'ほぼ常に空' が正常"},
    ]
    add_numbered_card_grid(slide, cards, n_cols=3,
                           x=ML, y=CY + 0.30, total_w=CW, total_h=AH - 0.40,
                           highlight_indices=list(range(len(cards))),
                           highlight_color=CORE_PURPLE,
                           font_name=FONT)
    set_footer(slide)


# ───────────────────────────────────────────────────────────────────────────
# Slide 6 — Processing Flow (Pattern P - Chevron Flow)
# ───────────────────────────────────────────────────────────────────────────
def slide_06_flow():
    slide = add_slide()
    add_breadcrumb(slide, "ServiceNow Research > アーキテクチャ")
    add_title(slide, "処理フロー")
    add_message_line(slide, "全工程で出典の公式性と URL 生存性を機械的に検証する")

    steps = [
        "① 質問受信",
        "② リリース判定",
        "③ 5並列リサーチ",
        "④ 集約 + URL検証",
        "⑤ ドラフト生成",
        "⑥ 客観レビュー",
        "⑦ 公開",
    ]
    flow_h = 0.85
    flow_y = CY + 0.30
    add_chevron_flow(slide, steps, x=ML, y=flow_y, total_w=CW, h=flow_h,
                     gap=0.10, fill_color=CORE_PURPLE, text_color=WHITE,
                     font_name=FONT, font_size_pt=12)

    # Detail text rows under flow
    detail_y = flow_y + flow_h + 0.25
    detail_h = AH - (flow_y - CY) - flow_h - 0.25 - 0.10
    details = [
        ("Step 1", "ユーザの自然文質問を受け付ける"),
        ("Step 2", "リリース語を検出、無ければ AskUserQuestion で確認"),
        ("Step 3", "5体の researcher を TeamCreate or 並列 Agent で同時起動"),
        ("Step 4", "allowlist チェック + HEAD で URL 生存確認、NotFound 排除"),
        ("Step 5", "脚注付き構造化 Markdown を生成 (結論/理由/具体例/出典)"),
        ("Step 6", "critique-reviewer が 6軸採点、< 0.85 なら 1回再生成"),
        ("Step 7", "ユーザへ最終回答を提示"),
    ]
    row_h = detail_h / len(details)
    for i, (k, v) in enumerate(details):
        y = detail_y + i * row_h
        # Number box
        nb = _add_rect(slide, ML, y, 1.0, row_h - 0.05, fill=CORE_PURPLE)
        tf = nb.text_frame
        tf.word_wrap = False
        _set_anchor_ctr(nb)
        _set_para(tf.paragraphs[0], k, size=12, bold=True, color=WHITE,
                  align=PP_ALIGN.CENTER, space_after_pt=0)
        # Description
        tb, tfd = _add_textbox(slide, ML + 1.10, y, CW - 1.10, row_h - 0.05,
                               anchor_ctr=True)
        _set_para(tfd.paragraphs[0], v, size=12, color=TEXT_BODY, space_after_pt=0)

    set_footer(slide)


# ───────────────────────────────────────────────────────────────────────────
# Slide 7 — Prerequisites (Pattern G - Table)
# ───────────────────────────────────────────────────────────────────────────
def slide_07_prereq():
    slide = add_slide()
    add_breadcrumb(slide, "ServiceNow Research > 導入")
    add_title(slide, "前提環境")
    add_message_line(slide, "stdlib のみで動作 — pip 依存ゼロで配布先の負担を最小化")

    # Build a table 5 rows x 4 cols
    headers = ["要件", "バージョン / 詳細", "確認コマンド", "必須度"]
    rows = [
        ["Claude Code",      "最新版",                  "claude --version",       "✅ 必須"],
        ["Python",           "3.10 以上",               "python3 --version",      "✅ 必須"],
        ["GitHub CLI (gh)",  "2.0+ / repo scope 必要",  "gh auth status",         "✅ 必須"],
        ["インターネット接続", "github.com / *.servicenow.com 到達可", "curl -I raw.githubusercontent.com", "✅ 必須"],
        ["pip パッケージ",   "PyYAML / requests など",  "(不要)",                "❌ 不要"],
    ]
    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl_y = CY + 0.40
    tbl_h = 4.20
    tbl = slide.shapes.add_table(n_rows, n_cols, Inches(ML), Inches(tbl_y),
                                 Inches(CW), Inches(tbl_h)).table
    # Column widths summing to CW (12.50)
    col_widths = [3.00, 3.80, 3.70, 2.00]
    for j, w in enumerate(col_widths):
        tbl.columns[j].width = Inches(w)

    # Header row
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = CORE_PURPLE
        cell.text_frame.clear()
        p = cell.text_frame.paragraphs[0]
        _set_para(p, h, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                  space_after_pt=0)

    # Body rows — single color (no alternating per brand rule)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = OFF_WHITE
            cell.text_frame.clear()
            p = cell.text_frame.paragraphs[0]
            color = TEXT_BODY
            if "❌" in val:
                color = MID_GRAY
            _set_para(p, val, size=13, color=color,
                      align=PP_ALIGN.CENTER if j >= 2 else PP_ALIGN.LEFT,
                      space_after_pt=0)

    set_footer(slide)


# ───────────────────────────────────────────────────────────────────────────
# Slide 8 — Setup Steps (Pattern X - Step Chart)
# ───────────────────────────────────────────────────────────────────────────
def slide_08_setup():
    slide = add_slide()
    add_breadcrumb(slide, "ServiceNow Research > 導入")
    add_title(slide, "導入手順 (3ステップ)")
    add_message_line(slide, "Zip 展開 → gh 認証 → Claude Code 再起動、所要 5分")

    phases = [
        {
            "label": "STEP 1",
            "color": CORE_PURPLE,
            "steps": [
                {
                    "title": "Zip を展開",
                    "subtitle": "配布物の配置",
                    "bullets": [
                        "$ unzip servicenow-research-1.0.0.zip \\",
                        "    -d ~/.claude/skills/",
                        "→ servicenow-research/ ディレクトリが展開される",
                        "",
                        "📌 ~/.claude/skills/ 直下が必須。サブディレクトリ不可",
                    ],
                },
            ],
        },
        {
            "label": "STEP 2",
            "color": CORE_PURPLE,
            "steps": [
                {
                    "title": "gh CLI を認証",
                    "subtitle": "未認証の場合のみ",
                    "bullets": [
                        "$ gh auth login",
                        "→ GitHub.com / HTTPS / Web ブラウザ を選択",
                        "→ 認証後 'gh auth status' で確認",
                        "",
                        "🔧 必要 scope: repo (標準ログインで自動付与)",
                    ],
                },
            ],
        },
        {
            "label": "STEP 3",
            "color": DARKEST_PURPLE,
            "steps": [
                {
                    "title": "Claude Code を再起動",
                    "subtitle": "スキル認識のため",
                    "bullets": [
                        "Claude Code を終了して立ち上げ直す",
                        "(/reload-plugins でも可)",
                        "",
                        "💡 動作確認: 自然文で「ServiceNow の Import Set について教えて」",
                        "→ AskUserQuestion でリリース確認が出れば成功",
                    ],
                },
            ],
        },
    ]
    add_step_chart(slide, phases, x=ML, y=CY + 0.30,
                   total_w=CW, total_h=AH - 0.40, font_name=FONT)
    set_footer(slide)


# ───────────────────────────────────────────────────────────────────────────
# Slide 9 — Usage Examples (Pattern Q - Icon Grid 6 items)
# ───────────────────────────────────────────────────────────────────────────
def slide_09_usage():
    slide = add_slide()
    add_breadcrumb(slide, "ServiceNow Research > 使い方")
    add_title(slide, "質問例 (自動発動)")
    add_message_line(slide, "ServiceNow と書かれた自然文質問で自動発動する")

    cells = [
        ("💡", "機能の概要", '"Now Assist for ITSM ってなに?"', "概念・機能説明"),
        ("🔧", "API リファレンス", '"GlideRecord の addQuery の使い方"', "API・メソッド調査"),
        ("📌", "設定手順", '"Yokohama で MID Server 構成手順"', "Step-by-step 取得"),
        ("🆕", "リリース差分", '"Australia と Zurich で Workflow Studio の違い"', "横断比較"),
        ("⚠️", "エラー対処", '"Service Portal で widget が表示されない原因"', "実務トラブル"),
        ("🔗", "比較・選定", '"Flow Designer と Workflow Studio の違い"', "機能比較"),
    ]
    n_cols, n_rows = 3, 2
    grid_y = CY + 0.30
    grid_h = AH - 0.40
    gap_h = 0.30
    gap_v = 0.30
    cell_w = (CW - (n_cols - 1) * gap_h) / n_cols
    cell_h = (grid_h - (n_rows - 1) * gap_v) / n_rows

    for idx, (icon, title_t, query, kind) in enumerate(cells):
        r, c = divmod(idx, n_cols)
        x = ML + c * (cell_w + gap_h)
        y = grid_y + r * (cell_h + gap_v)
        _add_rect(slide, x, y, cell_w, cell_h, fill=OFF_WHITE)
        # Icon (large)
        tb, tf = _add_textbox(slide, x + 0.20, y + 0.15, 0.80, 0.80, anchor_ctr=True)
        _set_para(tf.paragraphs[0], icon, size=32, color=CORE_PURPLE,
                  align=PP_ALIGN.CENTER, space_after_pt=0)
        # Kind tag
        tb, tf = _add_textbox(slide, x + 1.05, y + 0.20, cell_w - 1.20, 0.30)
        _set_para(tf.paragraphs[0], kind, size=11, color=CORE_PURPLE,
                  bold=True, space_after_pt=0)
        # Title
        tb, tf = _add_textbox(slide, x + 1.05, y + 0.50, cell_w - 1.20, 0.35)
        _set_para(tf.paragraphs[0], title_t, size=15, bold=True, color=BLACK,
                  space_after_pt=0)
        # Query example
        tb, tf = _add_textbox(slide, x + 0.20, y + cell_h - 0.85, cell_w - 0.40, 0.75,
                              anchor_ctr=True)
        _set_para(tf.paragraphs[0], query, size=12, color=TEXT_BODY, space_after_pt=0)

    set_footer(slide)


# ───────────────────────────────────────────────────────────────────────────
# Slide 10 — Release Coverage (Pattern Y - Arrow Roadmap)
# ───────────────────────────────────────────────────────────────────────────
def slide_10_release():
    slide = add_slide()
    add_breadcrumb(slide, "ServiceNow Research > 制約")
    add_title(slide, "対応リリース範囲")
    add_message_line(slide, "年2回の rolling 保持で常に最新4世代を網羅する")

    # Timeline arrow
    timeline_y = CY + 0.85
    timeline_h = 1.20
    n = 4
    arrow_w = (CW - (n - 1) * 0.20) / n
    releases = [
        ("Xanadu",    "2024 後半", "△ まもなく削除", LIGHT_GRAY,    BLACK),
        ("Yokohama",  "2025 前半", "○ 対応中",     LIGHTEST_PURPLE, BLACK),
        ("Zurich",    "2025 後半", "○ 対応中",     CORE_PURPLE,    WHITE),
        ("Australia", "2026 前半", "◎ 最新",       CORE_PURPLE,    WHITE),
    ]
    for i, (name, when, status, fill, fg) in enumerate(releases):
        x = ML + i * (arrow_w + 0.20)
        # Use HOME_PLATE shape for arrow effect
        try:
            shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.PENTAGON,
                                         Inches(x), Inches(timeline_y),
                                         Inches(arrow_w), Inches(timeline_h))
        except Exception:
            shp = _add_rect(slide, x, timeline_y, arrow_w, timeline_h, fill=fill)
        else:
            shp.fill.solid()
            shp.fill.fore_color.rgb = fill
            shp.line.fill.background()
        tf = shp.text_frame
        tf.word_wrap = True
        _set_anchor_ctr(shp)
        _set_para(tf.paragraphs[0], name, size=20, bold=True, color=fg,
                  align=PP_ALIGN.CENTER, space_after_pt=4)
        p2 = tf.add_paragraph()
        _set_para(p2, when, size=12, color=fg, align=PP_ALIGN.CENTER,
                  space_after_pt=4)
        p3 = tf.add_paragraph()
        _set_para(p3, status, size=14, bold=True, color=fg,
                  align=PP_ALIGN.CENTER, space_after_pt=0)

    # Below: deleted releases & maintenance note
    note_y = timeline_y + timeline_h + 0.40
    note_h = AH - (timeline_y - CY) - timeline_h - 0.40 - 0.10
    _add_rect(slide, ML, note_y, CW, note_h, fill=OFF_WHITE)
    tb, tf = _add_textbox(slide, ML + 0.30, note_y + 0.20, CW - 0.60, note_h - 0.40,
                          anchor_ctr=True)
    _multi_line(tf, [
        "📌 ローリング保持ポリシー — ServiceNow 公式が GA 時に最古ブランチを削除",
        "❌ Vancouver / Utah / Tokyo / Rome 等の旧リリース固有情報は対応外",
        "🔧 メンテ作業: 新リリース GA 時に scripts/raw_fetch.py の ALLOWED_BRANCHES を更新 (年2回)",
    ], size=15, color=TEXT_BODY, space_after_pt=10)

    set_footer(slide)


# ───────────────────────────────────────────────────────────────────────────
# Slide 11 — Output Sample (Pattern A - Title + Body, sole use)
# ───────────────────────────────────────────────────────────────────────────
def slide_11_sample():
    slide = add_slide()
    add_breadcrumb(slide, "ServiceNow Research > 使い方")
    add_title(slide, "出力サンプル")
    add_message_line(slide, "結論→理由→具体例→出典 の構造化、本文は脚注番号のみ")

    # Two-column layout: query on left, structured response on right
    col_gap = 0.40
    left_w = 3.50
    right_w = CW - left_w - col_gap
    body_y = CY + 0.30
    body_h = AH - 0.40

    # Left: query box
    _add_rect(slide, ML, body_y, left_w, body_h, fill=DARKEST_PURPLE)
    tb, tf = _add_textbox(slide, ML + 0.25, body_y + 0.30, left_w - 0.50, body_h - 0.60,
                          anchor_ctr=True)
    p = tf.paragraphs[0]
    _set_para(p, "💬 ユーザ質問", size=14, bold=True, color=WHITE, space_after_pt=12)
    p = tf.add_paragraph()
    _set_para(p, "ServiceNowのImportSetについて概要を教えて", size=18, bold=True,
              color=WHITE, space_after_pt=18)
    p = tf.add_paragraph()
    _set_para(p, "（リリース未指定 → AskUserQuestion で Australia を確認）",
              size=11, color=LIGHTEST_PURPLE, space_after_pt=12)
    p = tf.add_paragraph()
    _set_para(p, "↓ 5並列リサーチ + URL検証", size=12, color=WHITE,
              align=PP_ALIGN.CENTER, space_after_pt=0)

    # Right: response sample
    rx = ML + left_w + col_gap
    _add_rect(slide, rx, body_y, right_w, body_h, fill=OFF_WHITE)
    tb, tf = _add_textbox(slide, rx + 0.25, body_y + 0.20, right_w - 0.50, body_h - 0.40)
    sample = [
        ("## 結論",                                        16, True,  CORE_PURPLE),
        ("Import Set は ServiceNow の **データ取り込み staging** 機構で、",
                                                            13, False, TEXT_BODY),
        ("3段階パイプライン (Source → Staging → Target) [1]。", 13, False, TEXT_BODY),
        ("",                                                10, False, TEXT_BODY),
        ("## 理由 / 背景",                                  16, True,  CORE_PURPLE),
        ("外部データ構造と ServiceNow スキーマは一致しないため、",
                                                            13, False, TEXT_BODY),
        ("検証→変換→反映の安全な経路を提供する [1][2]。",   13, False, TEXT_BODY),
        ("",                                                10, False, TEXT_BODY),
        ("## 具体例 / 手順",                                16, True,  CORE_PURPLE),
        ("1. Data Sources > New で CSV をアップロード [3]",  13, False, TEXT_BODY),
        ("2. Transform Map で field 対応付け [2]",          13, False, TEXT_BODY),
        ("3. Coalesce キーで update vs insert を制御 [2]",  13, False, TEXT_BODY),
        ("",                                                10, False, TEXT_BODY),
        ("## 出典",                                         16, True,  CORE_PURPLE),
        ("1. [Imports overview](raw.../imports-overview-page.md)",
                                                            11, False, MID_GRAY),
        ("2. [Create a Transform Map](raw.../t_CreateATransformMap.md)",
                                                            11, False, MID_GRAY),
        ("3. [Getting Started Blog (developer)](.../blog.do?p=...)",
                                                            11, False, MID_GRAY),
    ]
    for j, (txt, sz, bold, col) in enumerate(sample):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        _set_para(p, txt, size=sz, bold=bold, color=col, space_after_pt=2)

    set_footer(slide)


# ───────────────────────────────────────────────────────────────────────────
# Slide 12 — Constraints (Pattern E - Bullet with GT Icon)
# ───────────────────────────────────────────────────────────────────────────
def slide_12_constraints():
    slide = add_slide()
    add_breadcrumb(slide, "ServiceNow Research > 制約")
    add_title(slide, "制約事項")
    add_message_line(slide, "「わかりません」が出るのは設計通りの正常動作")

    items = [
        ("🔒", "公式以外参照禁止",
         "Qiita / Stack Overflow / 個人ブログ / Reddit / YouTube は出典に含まれない。",
         "コードレベルで allowlist 強制 (Source constructor で ValueError)。"),
        ("🚫", "認証壁ソース対応外",
         "support.servicenow.com の KB 本文は SSO 必須のため取得不可。",
         "WebFetch が認証 redirect を返したら自動的に空ヒット扱い。"),
        ("📅", "古いリリース対応外",
         "Vancouver より古いリリースは ServiceNowDocs リポジトリから削除済。",
         "Australia / Zurich / Yokohama / Xanadu の最新4世代のみ対応。"),
        ("🛡️", "架空機能の捏造防止",
         "5ソース全てで該当無しなら render_unknown() で正直回答。",
         "「公式情報源では該当する情報が見つかりませんでした」+ 検索内訳開示。"),
    ]
    body_y = CY + 0.30
    body_h = AH - 0.40
    row_h = body_h / len(items)

    for i, (icon, head, b1, b2) in enumerate(items):
        y = body_y + i * row_h
        # Icon column (left)
        _add_rect(slide, ML, y, 1.10, row_h - 0.10, fill=CORE_PURPLE)
        tb, tf = _add_textbox(slide, ML, y, 1.10, row_h - 0.10, anchor_ctr=True)
        _set_para(tf.paragraphs[0], icon, size=28, color=WHITE, align=PP_ALIGN.CENTER,
                  space_after_pt=0)
        # Content column
        cx = ML + 1.30
        cw = CW - 1.30
        tb, tf = _add_textbox(slide, cx, y, cw, row_h - 0.10, anchor_ctr=True)
        p1 = tf.paragraphs[0]
        _set_para(p1, head, size=18, bold=True, color=BLACK, space_after_pt=4)
        p2 = tf.add_paragraph()
        _set_para(p2, b1, size=13, color=TEXT_BODY, space_after_pt=2)
        p3 = tf.add_paragraph()
        _set_para(p3, b2, size=13, color=TEXT_BODY, space_after_pt=0)

    set_footer(slide)


# ───────────────────────────────────────────────────────────────────────────
# Slide 13 — Troubleshooting (Pattern AB - Issue Tree)
# ───────────────────────────────────────────────────────────────────────────
def slide_13_troubleshoot():
    slide = add_slide()
    add_breadcrumb(slide, "ServiceNow Research > 制約")
    add_title(slide, "判別フロー — 「わかりません」頻発時")
    add_message_line(slide, "障害ではなく仕様による空回答かをまず切り分ける")

    body_y = CY + 0.20
    body_h = AH - 0.30

    # Root question
    root_w = 4.50
    root_h = 0.80
    root_x = ML + (CW - root_w) / 2
    root_y = body_y
    rb = _add_rect(slide, root_x, root_y, root_w, root_h, fill=DARKEST_PURPLE)
    tf = rb.text_frame
    tf.word_wrap = True
    _set_anchor_ctr(rb)
    _set_para(tf.paragraphs[0], "「わかりません」が頻発する?",
              size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER, space_after_pt=0)

    # Level 2: 4 branches
    n = 4
    branch_y = root_y + root_h + 0.45
    branch_h = 1.10
    branch_w = (CW - (n - 1) * 0.20) / n
    branches = [
        ("質問は ServiceNow 関連か?",  "No → スキル未発動",        "質問に 'ServiceNow' を明示", MID_GRAY),
        ("架空の機能名 / typo か?",     "Yes → 仕様通り正常",        "正しい機能名で再質問",      CORE_PURPLE),
        ("gh auth status は OK か?",   "No → 認証切れ",            "$ gh auth login を再実行", CORE_PURPLE),
        ("Vancouver 以前か?",           "Yes → 対応外",              "対応4世代から選び直す",     DARKEST_PURPLE),
    ]

    for i, (q, ans, action, color) in enumerate(branches):
        bx = ML + i * (branch_w + 0.20)
        # Question header
        _add_rect(slide, bx, branch_y, branch_w, 0.55, fill=color)
        tb, tf = _add_textbox(slide, bx + 0.10, branch_y, branch_w - 0.20, 0.55,
                              anchor_ctr=True)
        _set_para(tf.paragraphs[0], q, size=12, bold=True, color=WHITE,
                  align=PP_ALIGN.CENTER, space_after_pt=0)
        # Answer panel
        _add_rect(slide, bx, branch_y + 0.55, branch_w, branch_h, fill=OFF_WHITE)
        tb, tf = _add_textbox(slide, bx + 0.15, branch_y + 0.65,
                              branch_w - 0.30, branch_h - 0.20, anchor_ctr=True)
        p1 = tf.paragraphs[0]
        _set_para(p1, ans, size=12, bold=True, color=CORE_PURPLE, space_after_pt=4)
        p2 = tf.add_paragraph()
        _set_para(p2, "→ " + action, size=11, color=TEXT_BODY, space_after_pt=0)

        # Connector from root to branch
        bx_center = bx + branch_w / 2
        try:
            slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT,
                                       Inches(root_x + root_w / 2), Inches(root_y + root_h),
                                       Inches(bx_center), Inches(branch_y))
        except Exception:
            pass

    # Bottom: escape hatch
    bot_y = branch_y + 0.55 + branch_h + 0.30
    bot_h = body_h - (branch_y - body_y) - 0.55 - branch_h - 0.30
    if bot_h > 0.40:
        _add_rect(slide, ML, bot_y, CW, bot_h, fill=LIGHTEST_PURPLE)
        tb, tf = _add_textbox(slide, ML + 0.30, bot_y, CW - 0.60, bot_h,
                              anchor_ctr=True)
        _set_para(tf.paragraphs[0],
                  "💡 全てクリアでも頻発するなら → ネットワーク確認 (curl raw.githubusercontent.com) / scripts/url_check.py の TIMEOUT_SECONDS を社内プロキシ向けに延長",
                  size=12, color=TEXT_BODY, align=PP_ALIGN.CENTER, space_after_pt=0)

    set_footer(slide)


# ───────────────────────────────────────────────────────────────────────────
# Slide 14 — Summary + Distribution (Pattern J - KPI / Metrics)
# ───────────────────────────────────────────────────────────────────────────
def slide_14_summary():
    slide = add_slide()
    add_breadcrumb(slide, "ServiceNow Research > まとめ")
    add_title(slide, "まとめ — 配布物の概要")
    add_message_line(slide, "Zip を渡すだけで、5分で受領者が使い始められる")

    # KPI cards row
    kpi_y = CY + 0.30
    kpi_h = 2.70
    n = 4
    gap = 0.25
    card_w = (CW - (n - 1) * gap) / n
    kpis = [
        ("3.1 MB",      "配布物サイズ",        "Zip 圧縮後\nPPTX ガイド込み"),
        ("18 ファイル", "同梱内容",            "SKILL/scripts/refs/agents\n+ PPTX ガイド"),
        ("0",          "pip 依存",             "stdlib のみで動作\nセットアップ不要"),
        ("3 ステップ", "導入手順",             "Zip 展開→gh auth→再起動\n所要 5 分"),
    ]
    for i, (big, label, sub) in enumerate(kpis):
        x = ML + i * (card_w + gap)
        _add_rect(slide, x, kpi_y, card_w, kpi_h, fill=OFF_WHITE)
        # Big number
        tb, tf = _add_textbox(slide, x + 0.20, kpi_y + 0.15, card_w - 0.40, 1.00,
                              anchor_ctr=True)
        _set_para(tf.paragraphs[0], big, size=40, bold=True, color=CORE_PURPLE,
                  align=PP_ALIGN.CENTER, space_after_pt=0)
        # Label bar
        _add_rect(slide, x + 0.20, kpi_y + 1.20, card_w - 0.40, 0.55,
                  fill=CORE_PURPLE)
        tb, tf = _add_textbox(slide, x + 0.20, kpi_y + 1.20, card_w - 0.40, 0.55,
                              anchor_ctr=True, word_wrap=False)
        _set_para(tf.paragraphs[0], label, size=14, bold=True, color=WHITE,
                  align=PP_ALIGN.CENTER, space_after_pt=0)
        # Sub (taller for 2-line accommodation)
        tb, tf = _add_textbox(slide, x + 0.20, kpi_y + 1.80, card_w - 0.40, 0.85,
                              anchor_ctr=True)
        _multi_line(tf, sub.split("\n"), size=11, color=TEXT_BODY,
                    align=PP_ALIGN.CENTER, space_after_pt=2)

    # Bottom: distribution box
    box_y = kpi_y + kpi_h + 0.40
    box_h = AH - (kpi_y - CY) - kpi_h - 0.40 - 0.10
    _add_rect(slide, ML, box_y, CW, box_h, fill=DARKEST_PURPLE)
    tb, tf = _add_textbox(slide, ML + 0.40, box_y + 0.20, CW - 0.80, box_h - 0.40,
                          anchor_ctr=True)
    p1 = tf.paragraphs[0]
    _set_para(p1, "📦 配布方法 — Zip を共有相手に渡してください",
              size=18, bold=True, color=WHITE, space_after_pt=10)
    p2 = tf.add_paragraph()
    _set_para(p2, "$ unzip servicenow-research-1.0.0.zip -d ~/.claude/skills/",
              size=14, color=LIGHTEST_PURPLE, space_after_pt=4)
    p3 = tf.add_paragraph()
    _set_para(p3, "→ 詳細手順は同梱の README.md を参照", size=14,
              color=WHITE, space_after_pt=0)

    set_footer(slide)


# ───────────────────────────────────────────────────────────────────────────
# Build
# ───────────────────────────────────────────────────────────────────────────
def main():
    slide_01_cover()
    slide_02_agenda()
    slide_03_problem()
    slide_04_capabilities()
    slide_05_architecture()
    slide_06_flow()
    slide_07_prereq()
    slide_08_setup()
    slide_09_usage()
    slide_10_release()
    slide_11_sample()
    slide_12_constraints()
    slide_13_troubleshoot()
    slide_14_summary()

    make_closing_slide(prs, text_color=BLACK)
    strip_sections(prs)

    skill_root = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    out_path = os.path.join(skill_root, "servicenow-research_guide.pptx")
    prs.save(out_path)
    _sanitize_metadata(out_path)
    print(f"✅ Saved: {out_path}")


def _sanitize_metadata(pptx_path: str) -> None:
    """Replace template-inherited author identifiers with a generic label.

    python-pptx inherits docProps/core.xml + app.xml from the Accenture template,
    which embeds the original template authors' names. Without scrubbing, every
    regeneration would leak those identifiers. Idempotent: safe to re-run.
    """
    import re as _re
    import shutil as _shutil
    import tempfile as _tempfile
    import zipfile as _zipfile

    generic = "ServiceNow Research Skill"
    tmpdir = _tempfile.mkdtemp()
    extract_dir = os.path.join(tmpdir, "extract")
    os.makedirs(extract_dir)
    with _zipfile.ZipFile(pptx_path) as zf:
        zf.extractall(extract_dir)

    def _patch(path: str, replacements: dict) -> None:
        if not os.path.exists(path):
            return
        with open(path, encoding="utf-8") as f:
            s = f.read()
        for tag, new_val in replacements.items():
            s = _re.sub(
                rf"(<{tag}[^>]*>).*?(</{tag}>)",
                lambda m, v=new_val: m.group(1) + v + m.group(2),
                s, flags=_re.S,
            )
        with open(path, "w", encoding="utf-8") as f:
            f.write(s)

    _patch(os.path.join(extract_dir, "docProps", "core.xml"),
           {"dc:creator": generic, "cp:lastModifiedBy": generic})
    _patch(os.path.join(extract_dir, "docProps", "app.xml"),
           {"Author": generic, "LastAuthor": generic, "Company": "", "Manager": ""})

    new_pptx = pptx_path + ".tmp"
    with _zipfile.ZipFile(new_pptx, "w", _zipfile.ZIP_DEFLATED) as zf:
        for root_dir, _dirs, files in os.walk(extract_dir):
            for fname in files:
                full = os.path.join(root_dir, fname)
                arc = os.path.relpath(full, extract_dir)
                zf.write(full, arc)
    _shutil.move(new_pptx, pptx_path)
    _shutil.rmtree(tmpdir)


if __name__ == "__main__":
    main()
