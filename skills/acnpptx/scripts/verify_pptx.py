"""
PPTX Skill — Post-generation verification script.

Run after saving a PPTX to check for common issues:
    python scripts/verify_pptx.py output.pptx

Checks performed:
  1. Font size (nothing below 12pt, excluding footer/slidenum)
  2. Overflow (shapes within slide bounds)
  3. Overlap detection (significant shape overlaps)
  4. Footer verification (idx=20 + idx=21 on content slides)
  5. Text clipping detection (text overflows text box height)
  6. Horizontal text overflow (word_wrap=False width + textbox exceeds card boundary)
  7. Layout density (content fills 75%+ of available height)
  8. Per-slide audit summary (font range, shape count)

Approved fonts: Graphik, Meiryo UI, Arial, Meiryo, MS Gothic, MS PGothic,
                Yu Gothic, Yu Gothic UI, 游ゴシック, GT Sectra Fine
"""

import sys
import io

# Handle encoding on Windows terminals
if sys.platform == "win32":
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")

from pptx import Presentation
from pptx.util import Inches, Emu
from lxml import etree

_ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
NSMAP_ASVG = "http://schemas.microsoft.com/office/drawing/2016/SVG/main"
_footer_idxs = {20, 21}

# Approved font names (case-insensitive prefix match)
_APPROVED_FONTS = {
    "graphik", "meiryo ui", "meiryo", "arial",
    "ms gothic", "ms pgothic", "yu gothic", "yu gothic ui",
    "游ゴシック", "游明朝", "gt sectra fine",
    "calibri", "times new roman",  # system fallbacks
}


import re as _re

# Placeholder hint text patterns (case-insensitive)
_HINT_PATTERNS = [
    _re.compile(r"タイトルを入力", _re.IGNORECASE),
    _re.compile(r"テキストを入力", _re.IGNORECASE),
    _re.compile(r"サブタイトルを入力", _re.IGNORECASE),
    _re.compile(r"ここにテキスト", _re.IGNORECASE),
    _re.compile(r"マスター\s*テキストの書式", _re.IGNORECASE),
    _re.compile(r"マスター\s*タイトルの書式", _re.IGNORECASE),
    _re.compile(r"Place\s+subtitle\s+here", _re.IGNORECASE),
    _re.compile(r"Click\s+to\s+add", _re.IGNORECASE),
    _re.compile(r"Enter\s+text\s+here", _re.IGNORECASE),
    _re.compile(r"Add\s+title", _re.IGNORECASE),
    _re.compile(r"Add\s+text", _re.IGNORECASE),
    _re.compile(r"日付を入力", _re.IGNORECASE),
    _re.compile(r"フッター", _re.IGNORECASE),
]

# Exclude footer/page-number placeholders from hint check
_HINT_SKIP_IDXS = {20, 21}


def verify_pptx(filepath):
    """Comprehensive PPTX verification.

    Returns:
        tuple: (prs, errors, warnings, footer_slides, svg_total)
    """
    prs = Presentation(filepath)
    errors = []
    warnings = []

    # ── 0. Placeholder hint text residue check ────────────────
    for i, sl in enumerate(prs.slides, 1):
        for sh in sl.shapes:
            if sh.is_placeholder and sh.placeholder_format.idx in _HINT_SKIP_IDXS:
                continue
            if not sh.has_text_frame:
                continue
            text = sh.text_frame.text.strip()
            if not text:
                continue
            for pat in _HINT_PATTERNS:
                if pat.search(text):
                    snippet = text[:50] + ("..." if len(text) > 50 else "")
                    errors.append(
                        f"Slide {i}: '{sh.name}' にプレースホルダーのヒントテキストが残留 "
                        f"— 「{snippet}」"
                    )
                    break  # one match per shape is enough

    # ── 1. Font size check ───────────────────────────────────
    for i, sl in enumerate(prs.slides, 1):
        for sh in sl.shapes:
            if sh.is_placeholder and sh.placeholder_format.idx in _footer_idxs:
                continue
            for rp in sh._element.findall(".//a:rPr", _ns):
                sz = rp.get("sz")
                if sz and int(float(sz)) < 1200:
                    errors.append(
                        f"Slide {i}: font {int(float(sz))/100}pt < 12pt in '{sh.name}'"
                    )

    # ── 2. Overflow check ────────────────────────────────────
    sw = prs.slide_width
    sh_val = prs.slide_height
    tolerance = Emu(50000)  # ~0.05"
    for i, sl in enumerate(prs.slides, 1):
        for shp in sl.shapes:
            r = shp.left + shp.width
            b = shp.top + shp.height
            if r > sw + tolerance:
                errors.append(
                    f"Slide {i}: '{shp.name}' overflows RIGHT "
                    f"by {(r-sw)/914400:.2f}\""
                )
            if b > sh_val + tolerance:
                errors.append(
                    f"Slide {i}: '{shp.name}' overflows BOTTOM "
                    f"by {(b-sh_val)/914400:.2f}\""
                )

    # ── 3. Overlap detection ─────────────────────────────────
    margin = Emu(30000)
    for i, sl in enumerate(prs.slides, 1):
        shapes_bounds = []
        for shp in sl.shapes:
            if shp.width < Emu(50000) or shp.height < Emu(50000):
                continue
            shapes_bounds.append(
                (shp.name, shp.left, shp.top,
                 shp.left + shp.width, shp.top + shp.height)
            )
        for a_idx in range(len(shapes_bounds)):
            for b_idx in range(a_idx + 1, len(shapes_bounds)):
                a_name, a_l, a_t, a_r, a_b = shapes_bounds[a_idx]
                b_name, b_l, b_t, b_r, b_b = shapes_bounds[b_idx]
                a_contains_b = (
                    a_l - margin <= b_l and a_t - margin <= b_t
                    and a_r + margin >= b_r and a_b + margin >= b_b
                )
                b_contains_a = (
                    b_l - margin <= a_l and b_t - margin <= a_t
                    and b_r + margin >= a_r and b_b + margin >= a_b
                )
                if a_contains_b or b_contains_a:
                    continue
                ol = max(a_l, b_l)
                ot = max(a_t, b_t)
                orr = min(a_r, b_r)
                ob = min(a_b, b_b)
                if ol < orr and ot < ob:
                    ow = (orr - ol) / 914400
                    oh = (ob - ot) / 914400
                    area = ow * oh
                    a_area = (a_r - a_l) * (a_b - a_t) / (914400 * 914400)
                    b_area = (b_r - b_l) * (b_b - b_t) / (914400 * 914400)
                    min_area = min(a_area, b_area)
                    if min_area > 0 and area > 0.3 and area / min_area > 0.2:
                        warnings.append(
                            f"Slide {i}: '{a_name}' & '{b_name}' "
                            f"overlap {area:.1f} sq\""
                        )

    # ── 4. Footer verification ───────────────────────────────
    footer_slides = 0
    for i, sl in enumerate(prs.slides, 1):
        has_footer = any(
            ph.placeholder_format.idx == 20 for ph in sl.placeholders
        )
        has_num = any(
            ph.placeholder_format.idx == 21 for ph in sl.placeholders
        )
        if has_footer and has_num:
            footer_slides += 1

    # ── 5. Text clipping detection ───────────────────────────
    # Estimates required height including word-wrap line folding.
    # Single-paragraph textboxes are also checked (common clipping source).
    import math as _math

    def _get_font_size_emu(para):
        """Return font size in EMU for the first run in a paragraph, default 14pt."""
        for run in para.runs:
            if run.font.size:
                return int(run.font.size)
        for rp in para._p.findall(".//a:rPr", _ns):
            sz = rp.get("sz")
            if sz:
                return int(sz) * 127  # hundredths-of-pt → EMU
        return 14 * 12700  # 14pt default

    def _estimate_lines(text, fsz_emu, box_w_emu, word_wrap):
        """Estimate number of rendered lines accounting for word wrap."""
        if not text.strip():
            return 0
        if not word_wrap or box_w_emu <= 0:
            return 1
        # Approx char width: full-width CJK chars ≈ font_size, ASCII ≈ 0.55× font_size
        cjk_count = sum(1 for c in text if ord(c) > 0x3000)
        ascii_count = len(text) - cjk_count
        avg_char_w = (cjk_count * fsz_emu + ascii_count * fsz_emu * 0.55) / max(len(text), 1)
        if avg_char_w <= 0:
            return 1
        chars_per_line = max(1, int(box_w_emu / avg_char_w))
        return _math.ceil(len(text) / chars_per_line)

    for i, sl in enumerate(prs.slides, 1):
        for sh in sl.shapes:
            if sh.is_placeholder and sh.placeholder_format.idx in _footer_idxs:
                continue
            if not sh.has_text_frame:
                continue
            if sh.shape_type is not None and sh.shape_type == 19:  # table
                continue
            tf = sh.text_frame
            # Skip shapes with no visible text (decorative bars, empty containers)
            if not tf.text.strip():
                continue
            box_w_emu = sh.width
            word_wrap = tf.word_wrap  # True / False / None (None = default wrap)
            use_wrap = (word_wrap is None) or word_wrap

            # Estimate required height across all paragraphs
            est_h_emu = 0
            for p in tf.paragraphs:
                text = p.text
                fsz = _get_font_size_emu(p)
                line_h = int(fsz * 1.35)  # 1.35× line spacing
                n_lines = max(1, _estimate_lines(text, fsz, box_w_emu, use_wrap))
                space = int(p.space_before) if p.space_before else 0
                est_h_emu += line_h * n_lines + space

            shape_h = sh.height
            if est_h_emu > 0 and shape_h > 0 and est_h_emu > shape_h * 1.05:
                clip_pct = (est_h_emu - shape_h) / shape_h * 100
                level = "WARN" if clip_pct <= 20 else "CLIP"
                if clip_pct > 8:
                    msg = (
                        f"Slide {i}: '{sh.name}' text may be clipped "
                        f"(est={est_h_emu/914400:.2f}\" vs "
                        f"box={shape_h/914400:.2f}\", +{clip_pct:.0f}%)"
                    )
                    if level == "CLIP":
                        errors.append(msg)
                    else:
                        warnings.append(msg)

    # ── 6. Horizontal text overflow check ───────────────────
    # Detects two cases:
    #   (a) word_wrap=False textbox: estimated text width > box width
    #   (b) Any textbox whose RIGHT edge exceeds the right edge of a
    #       larger "container" shape that surrounds it (card overflow)
    for i, sl in enumerate(prs.slides, 1):
        # Build list of large background shapes (likely containers/cards)
        containers = []
        for sh in sl.shapes:
            if not sh.has_text_frame and sh.width >= Inches(0.80) and sh.height >= Inches(0.50):
                containers.append((sh.left, sh.top,
                                   sh.left + sh.width, sh.top + sh.height))

        for sh in sl.shapes:
            if sh.is_placeholder and sh.placeholder_format.idx in _footer_idxs:
                continue
            if not sh.has_text_frame:
                continue
            tf = sh.text_frame
            if not tf.text.strip():
                continue

            sh_r = sh.left + sh.width
            sh_b = sh.top + sh.height

            # (a) word_wrap=False: estimate text width
            if tf.word_wrap is False:
                for p in tf.paragraphs:
                    text = p.text
                    if not text.strip():
                        continue
                    fsz = _get_font_size_emu(p)
                    cjk = sum(1 for c in text if ord(c) > 0x3000)
                    asc = len(text) - cjk
                    est_w = int(cjk * fsz + asc * fsz * 0.60)
                    if est_w > sh.width * 1.10:
                        errors.append(
                            f"Slide {i}: '{sh.name}' text may overflow RIGHT "
                            f"(word_wrap=False, est={est_w/914400:.2f}\" "
                            f"vs box={sh.width/914400:.2f}\")"
                        )
                        break

            # (b) textbox right edge exceeds its surrounding container
            for c_l, c_t, c_r, c_b in containers:
                # Check if this textbox is positioned inside the container
                inside_x = c_l - Emu(36000) <= sh.left <= c_r
                inside_y = c_t - Emu(36000) <= sh.top <= c_b
                if inside_x and inside_y and sh_r > c_r + Emu(36000):
                    errors.append(
                        f"Slide {i}: '{sh.name}' right edge ({sh_r/914400:.2f}\") "
                        f"exceeds container right ({c_r/914400:.2f}\") "
                        f"by {(sh_r - c_r)/914400:.2f}\""
                    )
                    break  # report once per shape

    # ── 7. Layout density check ──────────────────────────────
    CY_emu = Inches(1.50)
    BY_emu = Inches(6.85)
    sparse_threshold = Inches(5.85)
    for i, sl in enumerate(prs.slides, 1):
        if i == 1:
            continue
        max_bottom = 0
        for sh in sl.shapes:
            if sh.is_placeholder and sh.placeholder_format.idx in _footer_idxs:
                continue
            bottom = sh.top + sh.height
            if sh.top >= CY_emu and bottom > max_bottom:
                max_bottom = bottom
        if max_bottom > 0 and max_bottom < sparse_threshold:
            fill_pct = (max_bottom - CY_emu) / (BY_emu - CY_emu) * 100
            warnings.append(
                f"Slide {i}: layout sparse -- content fills {fill_pct:.0f}% "
                f"of AH (bottom={max_bottom/914400:.2f}\", "
                f"target>={sparse_threshold/914400:.2f}\")"
            )

    # ── 8. Horizontal fill / centering check ────────────────
    # Content area should use most of the available width (CW).
    # Tables and charts should always be full-width.
    ML_emu = Inches(0.42)
    CW_emu = Inches(12.50)
    content_right_target = ML_emu + CW_emu  # expected right edge

    for i, sl in enumerate(prs.slides, 1):
        if i == 1 or i == len(prs.slides):
            continue  # skip cover/closing
        # Collect content-area shapes (below CY)
        content_shapes = []
        has_table_or_chart = False
        for sh in sl.shapes:
            if sh.is_placeholder and sh.placeholder_format.idx in _footer_idxs:
                continue
            if sh.top < CY_emu:
                continue
            content_shapes.append(sh)

        if not content_shapes:
            continue

        # Find bounding box of all content shapes
        min_left = min(sh.left for sh in content_shapes)
        max_right = max(sh.left + sh.width for sh in content_shapes)
        content_width = max_right - min_left

        # Check individual tables and charts
        for sh in content_shapes:
            is_tbl = sh.has_table if hasattr(sh, 'has_table') else False
            is_chart = sh.has_chart if hasattr(sh, 'has_chart') else False
            if is_tbl or is_chart:
                has_table_or_chart = True
                sh_w = sh.width
                kind = "テーブル" if is_tbl else "チャート"
                if sh_w < CW_emu * 0.80:
                    fill_pct = sh_w / CW_emu * 100
                    warnings.append(
                        f"Slide {i}: '{sh.name}' {kind}の幅が狭い "
                        f"({sh_w/914400:.1f}\" = CW の {fill_pct:.0f}%) — "
                        f"w=CW ({CW_emu/914400:.1f}\") でフル幅を使うこと"
                    )

        # Check overall content horizontal fill
        if content_width < CW_emu * 0.70:
            fill_pct = content_width / CW_emu * 100
            warnings.append(
                f"Slide {i}: コンテンツ幅が狭い "
                f"({content_width/914400:.1f}\" = CW の {fill_pct:.0f}%) — "
                f"右側に大きな空白。幅を拡大するかセンタリングすること"
            )

    # ── 9. Per-slide audit ───────────────────────────────────
    print("\n── Per-Slide Audit ──")
    svg_total = 0
    for i, sl in enumerate(prs.slides, 1):
        min_font, max_font = 9999, 0
        svg_n = 0
        for sh in sl.shapes:
            if sh.is_placeholder and sh.placeholder_format.idx in _footer_idxs:
                continue
            for rp in sh._element.findall(".//a:rPr", _ns):
                sz = rp.get("sz")
                if sz:
                    fsz = int(sz) / 100
                    min_font = min(min_font, fsz)
                    max_font = max(max_font, fsz)
            if sh._element.findall(
                ".//" + f"{{{NSMAP_ASVG}}}svgBlip"
            ):
                svg_n += 1
        svg_total += svg_n
        font_str = (
            f"fonts={min_font:.0f}-{max_font:.0f}pt"
            if max_font > 0
            else "no-text"
        )
        print(f"  Slide {i:2d}: {len(sl.shapes):3d} shapes, {font_str}")

    return prs, errors, warnings, footer_slides, svg_total


def main():
    if len(sys.argv) < 2:
        print("Usage: python verify_pptx.py <file.pptx>")
        sys.exit(1)

    filepath = sys.argv[1]
    prs, errors, warnings, footer_count, svg_total = verify_pptx(filepath)
    total_shapes = sum(len(sl.shapes) for sl in prs.slides)

    print(f"\nSlide count: {len(prs.slides)}")
    print(
        f"Slide dimensions: {prs.slide_width/914400:.3f}\" x "
        f"{prs.slide_height/914400:.3f}\""
    )
    print(f"Total shapes: {total_shapes}")
    print(f"SVG elements: {svg_total}")
    print(f"Footer/PageNum slides: {footer_count}/{len(prs.slides) - 1}")

    if errors:
        for e in errors:
            print(f"  FAIL: {e}")
    if warnings:
        for w in warnings:
            print(f"  WARN: {w}")
    if not errors and not warnings:
        print("\nALL CHECKS PASSED")
    elif not errors:
        print(f"\nWARNINGS ({len(warnings)}) — review above")
    else:
        print(f"\n{len(errors)} ERROR(S) FOUND — fix required")
        sys.exit(1)


if __name__ == "__main__":
    main()
