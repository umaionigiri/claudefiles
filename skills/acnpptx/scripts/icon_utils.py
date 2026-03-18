"""
PPTX Skill — Icon Library Utilities

Provides access to the Accenture icon library (Icon-library.pptx).

Two tiers of access:
  1. copy_icon_to_slide() — copies the icon's group shape directly from the
     source PPTX into the target slide (no extraction needed, always available)
  2. get_icon_path() — returns a pre-extracted PNG path if available
     (requires running extract_all_icons() once first)

Index file: ~/.claude/skills/acnpptx/assets/icon_index.json
  Format: { "keyword": { "slide": N, "shape": "shape_name", "png": "path_or_null" }, ... }

Usage:
    from icon_utils import add_icon, list_icons, find_icons

    # Add icon to slide at position (x, y) with given size in inches
    add_icon(slide, prs, "handshake", x=2.5, y=3.0, size=0.5)

    # Search for icons matching a keyword
    matches = find_icons("cloud")
"""

import os
import json
import copy

_SKILL_DIR  = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
# Icon library path — not bundled with skill; user must provide path to build_icon_index().
# Checked locations (in order): skill's assets/icons/, then None (trigger helpful error).
_ICON_LIB_CANDIDATES = [
    os.path.join(_SKILL_DIR, "assets", "icons", "Icon-library.pptx"),
    os.path.join(_SKILL_DIR, "assets", "Icon-library.pptx"),
]
_ICON_LIB = next((p for p in _ICON_LIB_CANDIDATES if os.path.isfile(p)), None)
_ICON_INDEX = os.path.join(_SKILL_DIR, "assets", "icon_index.json")
_ICONS_DIR  = os.path.join(_SKILL_DIR, "assets", "icons")


# ─────────────────────────────────────────────────────────────────────────────
# INDEX MANAGEMENT
# ─────────────────────────────────────────────────────────────────────────────

def build_icon_index(icon_lib_path=None, save=True):
    """
    Scan Icon-library.pptx and build a keyword index.

    Layout of Icon-library.pptx:
    - Each slide has icon shapes (Graphic/Freeform/Group) arranged in a grid
    - Below each icon column, a TextBox/Placeholder contains semicolon-separated
      keywords describing that icon (e.g., "cloud; download; storage")
    - Each keyword column has 4 icon variants (filled/outline × dark/light)

    Matching strategy:
    - For each keyword text box, find the nearest icon shape directly above it
      (same x-column, highest y that is still above the text box)
    - This ensures "cloud" text maps to the cloud icon, not a WiFi icon
      in the same column but further above

    Args:
        icon_lib_path: path to Icon-library.pptx (default: SkillSample location)
        save: write index to assets/icon_index.json

    Returns:
        dict: {keyword: {"slide": int, "shape": str, "png": None}}
    """
    import re
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    lib_path = icon_lib_path or _ICON_LIB
    if not lib_path or not os.path.isfile(lib_path):
        raise FileNotFoundError(
            "Icon-library.pptx not found. Place it at "
            f"{_ICON_LIB_CANDIDATES[0]} or pass the path explicitly: "
            "build_icon_index(icon_lib_path='/path/to/Icon-library.pptx')"
        )

    prs = Presentation(lib_path)
    index = {}
    icon_types = {MSO_SHAPE_TYPE.FREEFORM, MSO_SHAPE_TYPE.GROUP, MSO_SHAPE_TYPE.PICTURE}

    for slide_idx, slide in enumerate(prs.slides):
        text_boxes = []
        icon_shapes = []

        for sh in slide.shapes:
            x = (sh.left or 0) / 914400
            y = (sh.top or 0) / 914400
            w = (sh.width or 0) / 914400
            h = (sh.height or 0) / 914400

            if sh.has_text_frame:
                text = sh.text_frame.text.strip()
                if ";" in text and w < 3.0 and h < 1.5:
                    text_boxes.append({"x": x, "y": y, "w": w, "text": text})

            if sh.shape_type in icon_types and 0.1 < w < 2.0 and 0.1 < h < 2.0:
                icon_shapes.append({
                    "name": sh.name, "x": x, "y": y, "w": w, "h": h,
                    "cx": x + w / 2, "slide": slide_idx
                })

        # Match each text box to the nearest icon directly above it
        for tb in text_boxes:
            tx_left = tb["x"]
            tx_right = tb["x"] + tb["w"]
            tb_y = tb["y"]

            candidates = [
                ic for ic in icon_shapes
                if tx_left - 0.15 <= ic["cx"] <= tx_right + 0.15
                and ic["y"] < tb_y
            ]
            if not candidates:
                continue

            # Nearest icon above text (largest y = closest)
            candidates.sort(key=lambda ic: (-ic["y"], ic["x"]))
            best = candidates[0]

            keywords = [kw.strip().lower() for kw in re.split(r'[;,\n]', tb["text"]) if kw.strip()]
            for kw in keywords:
                kw = re.sub(r'[^\w\s-]', '', kw).strip()
                if len(kw) >= 2 and kw not in index:
                    index[kw] = {
                        "slide": best["slide"],
                        "shape": best["name"],
                        "png": None
                    }

    if save:
        os.makedirs(os.path.dirname(_ICON_INDEX), exist_ok=True)
        with open(_ICON_INDEX, "w", encoding="utf-8") as f:
            json.dump(index, f, ensure_ascii=False, indent=2)

    return index



def load_icon_index():
    """Load the icon index from disk. Returns empty dict if not found."""
    if not os.path.isfile(_ICON_INDEX):
        return {}
    with open(_ICON_INDEX, "r", encoding="utf-8") as f:
        return json.load(f)


def find_icons(keyword: str, max_results=10):
    """
    Search the icon index for shapes matching a keyword.

    Args:
        keyword: search term (case-insensitive, partial match)
        max_results: maximum number of results

    Returns:
        list of (keyword, info_dict) tuples sorted by relevance
    """
    index = load_icon_index()
    kw_lower = keyword.lower()
    exact, partial = [], []

    for k, info in index.items():
        if k == kw_lower:
            exact.append((k, info))
        elif kw_lower in k or k in kw_lower:
            partial.append((k, info))

    results = exact + partial
    return results[:max_results]


def list_icons(category_keyword=None):
    """
    List all icons in the index, optionally filtered by category keyword.

    Returns:
        sorted list of keyword strings
    """
    index = load_icon_index()
    keys = sorted(index.keys())
    if category_keyword:
        cat = category_keyword.lower()
        keys = [k for k in keys if cat in k]
    return keys


# ─────────────────────────────────────────────────────────────────────────────
# ICON EXTRACTION (one-time setup)
# ─────────────────────────────────────────────────────────────────────────────

def extract_all_icons(icon_lib_path=None, output_dir=None, use_com=True):
    """
    Extract all icons from Icon-library.pptx as PNG files.

    Requires pywin32 (PowerPoint COM) for best results on Windows.
    Falls back to pptx-based extraction for picture-type shapes.

    This is a one-time setup operation; icons are saved to assets/icons/.

    Args:
        icon_lib_path: path to Icon-library.pptx
        output_dir: where to save PNGs (default: assets/icons/)
        use_com: use PowerPoint COM renderer (Windows only)

    Returns:
        dict: {keyword: png_path}
    """
    lib_path = icon_lib_path or _ICON_LIB
    out_dir  = output_dir or _ICONS_DIR
    os.makedirs(out_dir, exist_ok=True)

    index = load_icon_index()
    if not index:
        index = build_icon_index(lib_path, save=False)

    extracted = {}

    if use_com and _com_available():
        extracted = _extract_via_com(lib_path, out_dir, index)
    else:
        extracted = _extract_via_pptx(lib_path, out_dir, index)

    # Update index with PNG paths
    for kw, png_path in extracted.items():
        if kw in index:
            index[kw]["png"] = png_path

    with open(_ICON_INDEX, "w", encoding="utf-8") as f:
        json.dump(index, f, ensure_ascii=False, indent=2)

    return extracted


def _com_available():
    try:
        import win32com.client
        return True
    except ImportError:
        return False


def _extract_via_com(lib_path, out_dir, index):
    """Extract icons by exporting each slide as PNG via PowerPoint COM."""
    import win32com.client
    import tempfile

    extracted = {}
    pptx_abs = os.path.abspath(lib_path)

    ppt = win32com.client.Dispatch("PowerPoint.Application")
    ppt.Visible = True

    try:
        pres = ppt.Presentations.Open(pptx_abs, ReadOnly=True, WithWindow=False)

        # Collect unique slides that have icons
        slides_needed = {}
        for kw, info in index.items():
            si = info["slide"] + 1  # COM is 1-indexed
            if si not in slides_needed:
                slides_needed[si] = []
            slides_needed[si].append((kw, info["shape"]))

        for slide_idx, kw_shapes in slides_needed.items():
            try:
                slide_com = pres.Slides(slide_idx)
                for kw, shape_name in kw_shapes:
                    # Find shape and export it
                    for shape_com in slide_com.Shapes:
                        if shape_com.Name == shape_name:
                            safe_name = "".join(
                                c if c.isalnum() or c in "-_" else "_"
                                for c in kw
                            )
                            png_path = os.path.join(out_dir, f"{safe_name}.png")
                            try:
                                shape_com.Export(os.path.abspath(png_path), "PNG")
                                extracted[kw] = png_path
                            except Exception:
                                pass
                            break
            except Exception:
                continue

        pres.Close()
    finally:
        ppt.Quit()

    return extracted


def _extract_via_pptx(lib_path, out_dir, index):
    """
    Extract picture-type icons via python-pptx (no COM needed).
    Only works for icons stored as embedded pictures (not grouped shapes).
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(lib_path)
    extracted = {}

    # Build slide → shapes map
    slides_map = {}
    for kw, info in index.items():
        si = info["slide"]
        if si not in slides_map:
            slides_map[si] = []
        slides_map[si].append((kw, info["shape"]))

    for slide_idx, kw_shapes in slides_map.items():
        slide = prs.slides[slide_idx]
        shape_lookup = {sh.name: sh for sh in slide.shapes}

        for kw, shape_name in kw_shapes:
            shape = shape_lookup.get(shape_name)
            if shape is None:
                continue
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                safe_name = "".join(
                    c if c.isalnum() or c in "-_" else "_"
                    for c in kw
                )
                png_path = os.path.join(out_dir, f"{safe_name}.png")
                try:
                    with open(png_path, "wb") as f:
                        f.write(shape.image.blob)
                    extracted[kw] = png_path
                except Exception:
                    pass

    return extracted


# ─────────────────────────────────────────────────────────────────────────────
# PLACING ICONS ON SLIDES
# ─────────────────────────────────────────────────────────────────────────────

def add_icon(slide, prs_target, keyword, x, y, size,
             icon_lib_path=None, color_override=None):
    """
    Add an icon to a slide.

    Lookup order:
      1. Pre-extracted PNG in assets/icons/  (if available)
      2. Copy group shape from Icon-library.pptx (always available)

    Args:
        slide       : target slide object
        prs_target  : Presentation object (needed for group copy)
        keyword     : icon name/keyword (e.g. "handshake", "cloud")
        x, y        : top-left position in inches
        size        : width and height in inches (square)
        icon_lib_path: path to Icon-library.pptx (default: auto-detected)
        color_override: not used for PNG mode; reserved for future group recolor

    Returns:
        shape object added to slide, or None if not found
    """
    from pptx.util import Inches

    # Try PNG first (pre-extracted)
    png_path = _get_png_path(keyword)
    if png_path and os.path.isfile(png_path):
        return _add_png_icon(slide, png_path, x, y, size)

    # Fall back to copying group shape from source PPTX
    lib_path = icon_lib_path or _ICON_LIB
    if not os.path.isfile(lib_path):
        print(f"[icon_utils] Icon library not found: {lib_path}")
        return None

    return _copy_icon_group(slide, lib_path, keyword, x, y, size)


def _get_png_path(keyword):
    """Return the PNG path for a keyword from the index, or None."""
    index = load_icon_index()
    kw_lower = keyword.lower()
    info = index.get(kw_lower)
    if info and info.get("png"):
        return info["png"]
    # Try partial match
    for k, v in index.items():
        if kw_lower in k and v.get("png"):
            return v["png"]
    return None


def _add_png_icon(slide, png_path, x, y, size):
    """Add a PNG icon to a slide."""
    from pptx.util import Inches
    from PIL import Image as PILImage

    img = PILImage.open(png_path)
    w, h = img.size
    ar = w / h
    if ar >= 1:
        iw, ih = size, size / ar
    else:
        iw, ih = size * ar, size
    # Center within the size×size box
    ox = x + (size - iw) / 2
    oy = y + (size - ih) / 2
    return slide.shapes.add_picture(png_path, Inches(ox), Inches(oy),
                                    Inches(iw), Inches(ih))


def _copy_icon_group(slide, lib_path, keyword, x, y, size):
    """
    Copy icon group shape from Icon-library.pptx into the target slide.

    The group is scaled and repositioned to (x, y, size×size).
    """
    from pptx import Presentation
    from pptx.util import Inches, Emu
    from lxml import etree

    index = load_icon_index()
    kw_lower = keyword.lower()

    # Find in index
    info = index.get(kw_lower)
    if not info:
        # Try partial
        for k, v in index.items():
            if kw_lower in k:
                info = v
                break
    if not info:
        print(f"[icon_utils] Icon '{keyword}' not found in index. "
              f"Run build_icon_index() first.")
        return None

    lib_prs = Presentation(lib_path)
    slide_idx = info["slide"]
    shape_name = info["shape"]

    if slide_idx >= len(lib_prs.slides):
        return None

    src_slide = lib_prs.slides[slide_idx]
    src_shape = None
    for sh in src_slide.shapes:
        if sh.name == shape_name:
            src_shape = sh
            break
    if src_shape is None:
        return None

    # Deep-copy the XML element into the target slide's spTree
    src_el = src_shape._element
    cloned = copy.deepcopy(src_el)

    # Reposition and resize
    target_emu = int(Inches(size))
    try:
        src_w = src_shape.width or Inches(0.5)
        src_h = src_shape.height or Inches(0.5)
        scale = min(target_emu / src_w, target_emu / src_h)

        new_w = int(src_w * scale)
        new_h = int(src_h * scale)
        new_x = int(Inches(x)) + (target_emu - new_w) // 2
        new_y = int(Inches(y)) + (target_emu - new_h) // 2

        # Find and update xfrm
        ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
        ns_p = "http://schemas.openxmlformats.org/presentationml/2006/main"

        for xfrm in cloned.iter("{%s}xfrm" % ns):
            for off in xfrm.findall("{%s}off" % ns):
                off.set("x", str(new_x))
                off.set("y", str(new_y))
            for ext in xfrm.findall("{%s}ext" % ns):
                ext.set("cx", str(new_w))
                ext.set("cy", str(new_h))
            break
    except Exception:
        pass

    slide.shapes._spTree.append(cloned)

    # Return a proxy (no clean python-pptx object for manually appended XML)
    return cloned


# ─────────────────────────────────────────────────────────────────────────────
# ICON GRID HELPER
# ─────────────────────────────────────────────────────────────────────────────

def add_icon_grid(slide, prs_target, items, x, y, total_w, total_h,
                  cols=3, icon_size=0.50, font_name="Meiryo UI",
                  font_size_pt=12, icon_lib_path=None):
    """
    Add a grid of icons with labels.

    Args:
        slide      : target slide
        prs_target : Presentation object
        items      : list of (keyword, label) tuples, e.g. [("cloud", "クラウド"), ...]
        x, y       : top-left of the grid in inches
        total_w    : total grid width in inches
        total_h    : total grid height in inches
        cols       : number of columns (default 3)
        icon_size  : icon square size in inches
        font_name  : label font
        font_size_pt: label font size

    Returns:
        list of added shapes
    """
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    n = len(items)
    rows = (n + cols - 1) // cols

    cell_w = total_w / cols
    cell_h = total_h / rows

    _TEXT_BODY = RGBColor(0x33, 0x33, 0x33)
    _MID_GRAY  = RGBColor(0x81, 0x81, 0x80)

    added = []
    for idx, (keyword, label) in enumerate(items):
        col = idx % cols
        row = idx // cols

        cx = x + col * cell_w
        cy = y + row * cell_h

        # Icon (centered in top part of cell)
        icon_x = cx + (cell_w - icon_size) / 2
        icon_y = cy + 0.05
        shape = add_icon(slide, prs_target, keyword, icon_x, icon_y, icon_size,
                         icon_lib_path=icon_lib_path)
        if shape:
            added.append(shape)

        # Label below icon
        label_y = icon_y + icon_size + 0.05
        label_h = cell_h - icon_size - 0.15
        tb = slide.shapes.add_textbox(
            Inches(cx), Inches(label_y),
            Inches(cell_w), Inches(max(label_h, 0.25))
        )
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        from pptx.enum.text import PP_ALIGN
        p.alignment = PP_ALIGN.CENTER
        p.text = label
        p.font.name = font_name
        p.font.size = Pt(font_size_pt)
        p.font.color.rgb = _TEXT_BODY
        added.append(tb)

    return added


# ─────────────────────────────────────────────────────────────────────────────
# CLI: one-time index build
# ─────────────────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    import sys
    lib = sys.argv[1] if len(sys.argv) > 1 else None
    print("Building icon index...")
    idx = build_icon_index(lib)
    print(f"  Indexed {len(idx)} keywords → {_ICON_INDEX}")
    print("Done.")
