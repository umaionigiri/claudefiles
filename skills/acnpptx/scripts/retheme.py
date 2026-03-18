"""
retheme.py — PPTX テーマ移行スクリプト

既存 PPTX のスライドコンテンツを保持したまま、
カラートークンとスライドマスターを新テーマに一括置き換える。

使い方 (bash):
    # ACN → Fiori
    PYTHONUTF8=1 python ~/.claude/skills/acnpptx/scripts/retheme.py deck.pptx fiori

    # Fiori → ACN（旧テーマ明示）
    PYTHONUTF8=1 python ~/.claude/skills/acnpptx/scripts/retheme.py deck.pptx accenture --from fiori --out deck_acn.pptx

Python API:
    from retheme import migrate
    migrate("deck.pptx", target_theme="fiori", source_theme="accenture")

処理内容:
  1. カラートークン対応表（旧 hex → 新 hex）を構築
  2. 新テンプレートを土台にしてスライドを移植
  3. schemeClr をソーステーマで解決して srgbClr に変換後、色を置換
  4. placeholder が新テンプレートに存在しない場合は free-floating textbox として移植
  5. 画像リレーションシップも移植（icon PNG 等）

制限事項:
  - チャート内部の色は置換されない（チャートはコンテンツとして保持されるが色はそのまま）
  - スライドマスター提供の装飾（ロゴ・GTシンボル・フッター）は新マスターに更新される
"""

import sys
import os
import copy
import json
import argparse
import glob
import colorsys

# Windows terminal encoding
if sys.platform == "win32":
    import io
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")

_SCRIPTS_DIR = os.path.dirname(os.path.abspath(__file__))
_SKILL_DIR   = os.path.dirname(_SCRIPTS_DIR)
_THEMES_DIR  = os.path.join(_SKILL_DIR, "assets", "themes")
_MASTERS_DIR = os.path.join(_SKILL_DIR, "assets", "masters")

sys.path.insert(0, _SCRIPTS_DIR)

from pptx import Presentation
from pptx.oxml.ns import qn
from lxml import etree as _et

_NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
_NS   = {"a": _NS_A}
_PH_SKIP = {20, 21}   # footer / page-num — 新マスターが自動提供
_SHAPE_TAGS = frozenset([
    qn("p:sp"), qn("p:pic"), qn("p:graphicFrame"),
    qn("p:grpSp"), qn("p:cxnSp"),
])


# ── ユーティリティ ─────────────────────────────────────────────────────────────

def _load_theme_data(theme_name: str) -> dict:
    path = os.path.join(_THEMES_DIR, f"{theme_name}.json")
    if not os.path.isfile(path):
        # name は大文字小文字を問わず検索
        for f in glob.glob(os.path.join(_THEMES_DIR, "*.json")):
            with open(f, encoding="utf-8") as fp:
                d = json.load(fp)
            if d.get("name", "").lower() == theme_name.lower():
                return d
        raise FileNotFoundError(f"Theme not found: {theme_name} (looked in {_THEMES_DIR})")
    with open(path, encoding="utf-8") as f:
        return json.load(f)


def _resolve_template_path(theme_data: dict) -> str:
    tpl = theme_data.get("template", "")
    if tpl:
        if not os.path.isabs(tpl):
            tpl = os.path.join(_SKILL_DIR, tpl)
        if os.path.isfile(tpl):
            return tpl
    # fallback: masters/<name>.pptx
    name = theme_data.get("name", "").lower()
    tpl = os.path.join(_MASTERS_DIR, f"{name}.pptx")
    if os.path.isfile(tpl):
        return tpl
    raise FileNotFoundError(
        f"Template PPTX not found for theme '{theme_data.get('name')}'. "
        f"Expected: {tpl}"
    )


def _build_color_map(src_tokens: dict, dst_tokens: dict) -> dict:
    """トークン名を介して 旧hex → 新hex のマップを構築。"""
    cmap = {}
    for token, src_hex in src_tokens.items():
        if token in dst_tokens:
            s = src_hex.lstrip("#").upper()
            d = dst_tokens[token].lstrip("#").upper()
            if s != d:
                cmap[s] = d
    return cmap


def _expand_token_colors(tokens: dict) -> dict:
    """
    tokens から中間色を計算して {alias: hex(6)} の拡張マップを返す。
    theme_selector.py の _mid / _mid2 と同じ計算で
    LIGHT_PURPLE・DARK_PURPLE 相当の中点色を生成する。
    """
    def _mid(h1, h2):
        h1, h2 = h1.lstrip("#"), h2.lstrip("#")
        r = (int(h1[0:2], 16) + int(h2[0:2], 16)) // 2
        g = (int(h1[2:4], 16) + int(h2[2:4], 16)) // 2
        b = (int(h1[4:6], 16) + int(h2[4:6], 16)) // 2
        return f"{r:02X}{g:02X}{b:02X}"

    expanded = {k: v.lstrip("#").upper() for k, v in tokens.items()}
    if "primary" in tokens and "primary_dark" in tokens:
        expanded["_mid_pd"] = _mid(tokens["primary"], tokens["primary_dark"])
    if "primary" in tokens and "primary_light" in tokens:
        expanded["_mid_pl"] = _mid(tokens["primary"], tokens["primary_light"])
    if "primary_dark" in tokens and "primary_light" in tokens:
        expanded["_mid_dl"] = _mid(tokens["primary_dark"], tokens["primary_light"])
    return expanded


def _collect_slide_colors(src_prs) -> set:
    """src PPTX の全スライドから実際に使われている srgbClr の hex set を返す。"""
    colors = set()
    for slide in src_prs.slides:
        for el in slide._element.iter(f"{{{_NS_A}}}srgbClr"):
            val = el.get("val", "").upper()
            if len(val) == 6:
                colors.add(val)
    return colors


def _saturation(hex_color: str) -> float:
    """RGB hex から HSL 彩度 (0.0–1.0) を返す。"""
    r = int(hex_color[0:2], 16) / 255
    g = int(hex_color[2:4], 16) / 255
    b = int(hex_color[4:6], 16) / 255
    _, _, s = colorsys.rgb_to_hls(r, g, b)
    return s


def _build_color_map_fuzzy(
    src_tokens: dict, dst_tokens: dict,
    actual_colors: set, delta: int = 25,
    direct_delta: int = 442,
) -> dict:
    """
    color_map を3段階で構築する:

      1. token 名 exact match（primary → primary 等）
      2. src 中間色を含む近似マッチ（スライドの実色が src_tokens の中間色に近い場合）
      3. 汎用直接マッチ: actual_colors の有彩色を dst_tokens の最近傍有彩色に変換
         src_tokens が不明・旧バージョンでも機能する汎用アプローチ。
         ニュートラル色（彩度 < 0.15）は対象外。
         direct_delta のデフォルト値は RGB 空間の最大距離（≈442）のため、
         有彩色は必ず最近傍の dst 有彩色に変換される。

    Args:
        src_tokens:   ソーステーマの tokens dict（空でも可）
        dst_tokens:   移行先テーマの tokens dict
        actual_colors: スライド内の実 srgbClr hex set（大文字 6 桁）
        delta:        Phase 2 の最大 RGB ユークリッド距離（デフォルト 25）
        direct_delta: Phase 3 の最大 RGB ユークリッド距離（デフォルト 442 = 上限なし）
    """
    _SAT_THRESHOLD = 0.15  # 彩度この値未満はニュートラル色（黒・白・グレー）として変換スキップ

    def _dist(h1: str, h2: str) -> float:
        r1, g1, b1 = int(h1[0:2], 16), int(h1[2:4], 16), int(h1[4:6], 16)
        r2, g2, b2 = int(h2[0:2], 16), int(h2[2:4], 16), int(h2[4:6], 16)
        return ((r1 - r2) ** 2 + (g1 - g2) ** 2 + (b1 - b2) ** 2) ** 0.5

    # ── Phase 1: token 名 exact match ──────────────────────────────────────
    cmap = _build_color_map(src_tokens, dst_tokens)

    if not actual_colors:
        return cmap

    # ── Phase 2: src 中間色含む近似マッチ ──────────────────────────────────
    if src_tokens:
        src_exp = _expand_token_colors(src_tokens)
        dst_exp = _expand_token_colors(dst_tokens)
        src_to_dst: dict[str, str] = {}
        for key, src_hex in src_exp.items():
            if key in dst_exp and src_hex != dst_exp[key]:
                src_to_dst[src_hex] = dst_exp[key]

        if src_to_dst:
            src_keys = list(src_to_dst.keys())
            fuzzy_added = 0
            for actual in actual_colors:
                if actual in cmap:
                    continue
                best_key, best_d = None, float("inf")
                for sk in src_keys:
                    d = _dist(actual, sk)
                    if d < best_d:
                        best_d, best_key = d, sk
                if best_key is not None and best_d <= delta:
                    cmap[actual] = src_to_dst[best_key]
                    fuzzy_added += 1
            if fuzzy_added:
                print(f"  Fuzzy color matches (src-based): {fuzzy_added} color(s)")

    # ── Phase 3: 汎用直接マッチ（どんなソーステーマでも機能）──────────────
    if dst_tokens:
        dst_exp = _expand_token_colors(dst_tokens)
        # dst の有彩色リスト（ニュートラル色は変換先候補から除外）
        dst_chromatic_vals = [v for v in dst_exp.values() if _saturation(v) >= _SAT_THRESHOLD]
        dst_hex_set = set(dst_chromatic_vals)  # 既に dst テーマの色 → 変換不要

        direct_added = 0
        for actual in actual_colors:
            if actual in cmap:
                continue  # Phase 1/2 で既に変換対象
            if actual in dst_hex_set:
                continue  # 既に dst テーマの色なので変換不要
            if _saturation(actual) < _SAT_THRESHOLD:
                continue  # 黒・白・グレー等のニュートラル色はスキップ

            best_dst, best_d = None, float("inf")
            for dst_hex in dst_chromatic_vals:
                d = _dist(actual, dst_hex)
                if d < best_d:
                    best_d, best_dst = d, dst_hex
            if best_dst is not None and best_d <= direct_delta:
                cmap[actual] = best_dst
                direct_added += 1

        if direct_added:
            print(f"  Direct chromatic matches: {direct_added} color(s)")

    return cmap


def _replace_colors(element, color_map: dict) -> int:
    """XML ツリー内の全 srgbClr を color_map に従って置換。置換件数を返す。"""
    count = 0
    for el in element.iter(f"{{{_NS_A}}}srgbClr"):
        val = el.get("val", "").upper()
        if val in color_map:
            el.set("val", color_map[val])
            count += 1
    return count


# ── レイアウトタイプ推定 ────────────────────────────────────────────────────────

_COVER_KEYWORDS   = {"cover", "title slide", "タイトル スライド", "タイトルスライド", "表紙"}
_SECTION_KEYWORDS = {"section", "セクション", "divider", "chapter", "区切り", "section header"}
_BLANK_KEYWORDS   = {"blank", "空白"}


def _classify_layout_type(layout_name: str) -> str:
    """レイアウト名からスライドタイプ (cover/section/content/blank) を推定する。"""
    name = layout_name.lower()
    for kw in _COVER_KEYWORDS:
        if kw in name:
            return "cover"
    for kw in _SECTION_KEYWORDS:
        if kw in name:
            return "section"
    for kw in _BLANK_KEYWORDS:
        if kw in name:
            return "blank"
    return "content"


def _build_dst_layout_map(dst_prs) -> dict:
    """
    移行先テンプレートのレイアウトをプレースホルダー構造で判定して
    type → layout でマップする。レイアウト名に依存しないため異なるテンプレートでも動作する。

    判定ロジック:
      cover   : CENTER_TITLE (type=3) + SUBTITLE (type=4) の両方を持つ
      section : CENTER_TITLE (type=3) を持ち SUBTITLE (type=4) を持たない
      content : TITLE (type=1) + BODY (type=2) + footer idx (20 or 21) を持つ
      blank   : content と同じレイアウトにフォールバック
    """
    # placeholder type 定数（pp_placeholder_type の整数値）
    _CT  = 3   # CENTER_TITLE
    _T   = 1   # TITLE
    _ST  = 4   # SUBTITLE
    _B   = 2   # BODY

    layout_map: dict = {}

    for layout in dst_prs.slide_layouts:
        ph_types = {ph.placeholder_format.type for ph in layout.placeholders}
        ph_idxs  = {ph.placeholder_format.idx  for ph in layout.placeholders}

        # cover: CENTER_TITLE + SUBTITLE
        if "cover" not in layout_map:
            if _CT in ph_types and _ST in ph_types:
                layout_map["cover"] = layout

        # section: CENTER_TITLE, no SUBTITLE
        if "section" not in layout_map:
            if _CT in ph_types and _ST not in ph_types:
                layout_map["section"] = layout

        # content: TITLE (not CENTER_TITLE), BODY, footer idx=20 or 21
        if "content" not in layout_map:
            if (_T in ph_types and _CT not in ph_types and
                    _B in ph_types and (20 in ph_idxs or 21 in ph_idxs)):
                layout_map["content"] = layout

        # content_fallback1: TITLE + BODY (footer なし)
        if "content" not in layout_map and "_content_fb1" not in layout_map:
            if _T in ph_types and _CT not in ph_types and _B in ph_types:
                layout_map["_content_fb1"] = layout

        # content_fallback2: TITLE のみ（最もシンプルなレイアウト＝"Title Only" 相当）
        # ph数が最少のものを優先（装飾要素が少ない）
        if "content" not in layout_map:
            if (_T in ph_types and _CT not in ph_types and
                    _B not in ph_types and _ST not in ph_types):
                prev = layout_map.get("_content_fb2")
                if prev is None or len(list(layout.placeholders)) < len(list(prev.placeholders)):
                    layout_map["_content_fb2"] = layout

    # content の段階的フォールバック
    # fb2 (Title Only) が最もクリーンなキャンバス → fb1 (Title+Body) より優先
    if "content" not in layout_map:
        if "_content_fb2" in layout_map:
            layout_map["content"] = layout_map["_content_fb2"]
        elif "_content_fb1" in layout_map:
            layout_map["content"] = layout_map["_content_fb1"]
    # 一時キーを除去
    layout_map.pop("_content_fb1", None)
    layout_map.pop("_content_fb2", None)

    # fallback by index
    n = len(dst_prs.slide_layouts)
    if "cover" not in layout_map:
        layout_map["cover"] = dst_prs.slide_layouts[0]
    if "content" not in layout_map:
        # 最終フォールバック: idx=1 ("Title Only" が多いテンプレートで安全)
        layout_map["content"] = dst_prs.slide_layouts[min(1, n - 1)]
    if "section" not in layout_map:
        layout_map["section"] = dst_prs.slide_layouts[min(6, n - 1)]
    layout_map["blank"] = layout_map.get("content", dst_prs.slide_layouts[0])

    print("Destination layout map:")
    for lt, layout in layout_map.items():
        print(f"  {lt:10s} → [{layout.name}]")
    print()

    return layout_map


def _remove_empty_placeholders(slide):
    """
    コンテンツのない（空の）テキストプレースホルダーをスライド XML から削除する。
    これにより「タイトルを入力」「テキストを入力」などのヒントテキストが
    レンダリングに表示されるのを防ぐ。
    対象: テキストフレームを持ち、実テキストが空のプレースホルダー。
    除外: footer/pagenum (idx 20/21)、テキストフレームを持たないもの。
    """
    sp_tree = slide.shapes._spTree
    for el in list(sp_tree):
        if not _is_placeholder(el):
            continue
        idx = _get_ph_idx(el)
        if idx in _PH_SKIP:
            continue
        tx_body = el.find(f".//{qn('p:txBody')}")
        if tx_body is None:
            continue  # 画像プレースホルダー等はスキップ
        all_text = "".join(
            (r.text or "") for r in tx_body.iter(f"{{{_NS_A}}}r")
        )
        if not all_text.strip():
            sp_tree.remove(el)


def _is_placeholder(el) -> bool:
    return el.find(f".//{qn('p:ph')}") is not None


def _get_ph_idx(el) -> int | None:
    """p:ph 要素の idx を返す。なければ None。"""
    ph = el.find(f".//{qn('p:ph')}")
    if ph is None:
        return None
    return int(ph.get("idx", "0"))


# ── schemeClr 解決 ─────────────────────────────────────────────────────────────

def _build_src_theme_map(src_prs) -> dict:
    """
    ソース PPTX のマスターテーマ XML から schemeClr スロット名 → hex を構築。
    master_to_theme._build_theme_map と同ロジック。
    """
    theme_map: dict = {}
    for master in src_prs.slide_masters:
        for rel in master.part.rels.values():
            if "theme" not in rel.reltype.lower():
                continue
            tree = _et.fromstring(rel.target_part.blob)
            clr_el = tree.find(".//a:clrScheme", _NS)
            if clr_el is None:
                continue
            for child in clr_el:
                tag = child.tag.split("}")[1]
                for sub in child:
                    val = sub.get("val") or sub.get("lastClr", "")
                    if not val:
                        continue
                    if val.upper() in ("WINDOWTEXT", "WINDOW"):
                        theme_map[tag] = "#000000" if "dk" in tag.lower() or tag == "tx1" else "#FFFFFF"
                    else:
                        theme_map[tag] = "#" + val.upper()
            break
    if "bg1" not in theme_map:
        theme_map["bg1"] = theme_map.get("lt1", "#FFFFFF")
    if "tx1" not in theme_map:
        theme_map["tx1"] = theme_map.get("dk1", "#000000")
    return theme_map


def _resolve_scheme_colors(element, src_theme_map: dict, color_map: dict):
    """
    XML ツリー内の全 schemeClr を src_theme_map で hex に解決して srgbClr に変換し、
    color_map も適用する。solidFill だけでなく gradFill/gs・tcPr なども対象にする。
    """
    _FILL_CONTAINERS = {
        f"{{{_NS_A}}}solidFill",
        f"{{{_NS_A}}}gs",        # gradFill/gsLst/gs（グラデーションストップ）
    }
    for sf in element.iter():
        if sf.tag not in _FILL_CONTAINERS:
            continue
        for child in list(sf):
            tag = child.tag.split("}")[1]
            if tag != "schemeClr":
                continue
            base = src_theme_map.get(child.get("val", ""))
            if not base:
                continue

            # lumMod / lumOff / shade / tint を HLS 空間で適用
            r0 = int(base[1:3], 16)
            g0 = int(base[3:5], 16)
            b0 = int(base[5:7], 16)
            h, l, s = colorsys.rgb_to_hls(r0 / 255, g0 / 255, b0 / 255)
            for mod in child:
                mod_tag = mod.tag.split("}")[1]
                v = int(mod.get("val", "100000")) / 100000
                if mod_tag == "lumMod":
                    l = max(0.0, min(1.0, l * v))
                elif mod_tag == "lumOff":
                    l = max(0.0, min(1.0, l + v))
                elif mod_tag == "shade":
                    l = max(0.0, min(1.0, l * v))
                elif mod_tag == "tint":
                    l = max(0.0, min(1.0, l + (1 - l) * v))
            r1, g1, b1 = colorsys.hls_to_rgb(h, l, s)
            resolved = f"{int(r1 * 255):02X}{int(g1 * 255):02X}{int(b1 * 255):02X}"

            # color_map で置換
            final = color_map.get(resolved, resolved)

            # schemeClr → srgbClr に差し替え
            sf.remove(child)
            srgb = _et.SubElement(sf, f"{{{_NS_A}}}srgbClr")
            srgb.set("val", final)


# ── スライドコピー ─────────────────────────────────────────────────────────────

def _copy_image_rel(src_part, dst_part, old_rId: str) -> str | None:
    """src_part の画像を dst_part にコピーして新 rId を返す。"""
    try:
        rel = src_part.rels.get(old_rId)
        if rel is None:
            return None
        return dst_part.relate_to(rel.target_part, rel.reltype)
    except Exception as e:
        print(f"    Warning: image copy failed ({old_rId}): {e}")
        return None


def _fix_pic_rels(el_copy, src_slide, dst_slide):
    """
    el_copy ツリー内の全 p:pic 要素の r:embed rId を src → dst に移植する。
    p:pic 直下・グループ内・プレースホルダー変換後の全パターンに対応。
    """
    for pic in el_copy.iter(qn("p:pic")):
        blip = pic.find(f".//{{{_NS_A}}}blip")
        if blip is None:
            continue
        old_rId = blip.get(qn("r:embed"))
        if not old_rId:
            continue
        new_rId = _copy_image_rel(src_slide.part, dst_slide.part, old_rId)
        if new_rId:
            blip.set(qn("r:embed"), new_rId)


def _copy_shape_element(el, src_slide, dst_slide, dst_tree):
    """
    1つの図形要素を dst_tree に追加する（画像/チャートの rId 移植含む）。
    p:pic (直接) と p:grpSp 内のネストした p:pic も _fix_pic_rels で処理する。
    """
    el_copy = copy.deepcopy(el)

    # 画像 (p:pic) / グループ図形 (p:grpSp) — 全 p:pic の r:embed を移植
    if el.tag in (qn("p:pic"), qn("p:grpSp")):
        _fix_pic_rels(el_copy, src_slide, dst_slide)

    # graphicFrame (チャート・表) — rId 移植
    elif el.tag == qn("p:graphicFrame"):
        for rel_el in el.iter():
            chart_rId = rel_el.get(qn("r:id"))
            if chart_rId and chart_rId.startswith("rId"):
                rel = src_slide.part.rels.get(chart_rId)
                if rel:
                    try:
                        new_rId = dst_slide.part.relate_to(rel.target_part, rel.reltype)
                        el_copy_rel_el = el_copy.find(
                            ".//" + rel_el.tag
                        ) if rel_el is not el else el_copy
                        if el_copy_rel_el is not None:
                            el_copy_rel_el.set(qn("r:id"), new_rId)
                    except Exception:
                        pass

    dst_tree.append(el_copy)


def _inject_xfrm(el_copy, sh):
    """
    python-pptx shape の解決済み座標を el_copy の p:spPr/a:xfrm に明示的に埋め込む。
    位置がレイアウト継承のみで slide XML に記録されていない場合に必要。
    """
    if sh.left is None or sh.top is None or sh.width is None or sh.height is None:
        return
    spPr = el_copy.find(qn("p:spPr"))
    if spPr is None:
        return
    # 既存の xfrm を削除
    existing = spPr.find(f"{{{_NS_A}}}xfrm")
    if existing is not None:
        spPr.remove(existing)
    # 新規 xfrm を spPr の先頭に挿入
    xfrm = _et.Element(f"{{{_NS_A}}}xfrm")
    off  = _et.SubElement(xfrm, f"{{{_NS_A}}}off")
    off.set("x", str(sh.left))
    off.set("y", str(sh.top))
    ext  = _et.SubElement(xfrm, f"{{{_NS_A}}}ext")
    ext.set("cx", str(sh.width))
    ext.set("cy", str(sh.height))
    spPr.insert(0, xfrm)


def _copy_shapes(src_slide, dst_slide, dst_ph_idxs: set):
    """
    非プレースホルダー図形を src から dst にコピー。
    さらに、dst に存在しない idx のプレースホルダーは p:ph を除去して
    free-floating textbox として移植する（カバー・メッセージラインの消失防止）。
    座標がレイアウト継承の場合は python-pptx の解決済み値を明示的に埋め込む。
    """
    src_tree = src_slide.shapes._spTree
    dst_tree = dst_slide.shapes._spTree

    # python-pptx shape を idx → shape でマップ（位置解決用）
    src_ph_map = {
        sh.placeholder_format.idx: sh
        for sh in src_slide.shapes
        if sh.is_placeholder
    }

    # dst のデフォルト非プレースホルダー図形を削除
    for el in list(dst_tree):
        if el.tag in _SHAPE_TAGS and not _is_placeholder(el):
            dst_tree.remove(el)

    for el in src_tree:
        if el.tag not in _SHAPE_TAGS:
            continue

        if _is_placeholder(el):
            idx = _get_ph_idx(el)
            # footer/page-num はスキップ
            if idx in _PH_SKIP:
                continue
            # dst にマッチする idx がある → _copy_placeholder_text が処理するのでスキップ
            if idx in dst_ph_idxs:
                continue
            # dst に存在しない idx → p:ph を外して free-floating textbox として移植
            el_copy = copy.deepcopy(el)
            for ph_el in el_copy.findall(f".//{qn('p:ph')}"):
                ph_el.getparent().remove(ph_el)
            # 座標がレイアウト継承のみの場合、解決済み座標を明示埋め込み
            if idx in src_ph_map:
                _inject_xfrm(el_copy, src_ph_map[idx])
            # ピクチャープレースホルダー (p:pic) の画像 rId を移植
            _fix_pic_rels(el_copy, src_slide, dst_slide)
            dst_tree.append(el_copy)
        else:
            _copy_shape_element(el, src_slide, dst_slide, dst_tree)


def _copy_placeholder_text(src_slide, dst_slide, dst_ph_idxs: set):
    """
    dst に存在する idx のプレースホルダーのテキストを src から dst にコピー。
    """
    for src_ph in src_slide.placeholders:
        idx = src_ph.placeholder_format.idx
        if idx in _PH_SKIP:
            continue
        if idx not in dst_ph_idxs:
            continue
        try:
            dst_ph = dst_slide.placeholders[idx]
        except KeyError:
            continue
        if src_ph.has_text_frame and dst_ph.has_text_frame:
            src_body = src_ph.text_frame._txBody
            dst_body = dst_ph.text_frame._txBody
            dst_body.getparent().replace(dst_body, copy.deepcopy(src_body))


# ── 自動テーマ検出 ─────────────────────────────────────────────────────────────

def _auto_detect_source_theme(src_path: str) -> dict:
    """
    利用可能な全テーマとの色一致数で旧テーマを自動検出し tokens を返す。
    一致数が0の場合は空辞書を返す（色置換なしで続行）。
    """
    theme_files = [
        f for f in glob.glob(os.path.join(_THEMES_DIR, "*.json"))
        if not os.path.basename(f).startswith("_")
    ]
    if not theme_files:
        print("Warning: no theme files found, color replacement skipped")
        return {}

    from master_to_theme import _collect_colors_from_prs
    src_prs = Presentation(src_path)
    fill_c, text_c, _ = _collect_colors_from_prs(src_prs)
    src_colors = {c.lstrip("#").upper() for c in list(fill_c) + list(text_c)}

    best_name, best_tokens, best_score = None, {}, -1
    for tf in theme_files:
        with open(tf, encoding="utf-8") as f:
            data = json.load(f)
        tokens = data.get("tokens", {})
        score = sum(
            1 for v in tokens.values()
            if v.lstrip("#").upper() in src_colors
        )
        if score > best_score:
            best_score, best_name, best_tokens = score, data.get("name"), tokens

    if best_score > 0:
        print(f"Auto-detected source theme: {best_name} ({best_score} color matches)")
        return best_tokens

    print("Warning: could not auto-detect source theme, color replacement skipped")
    return {}


# ── メイン処理 ────────────────────────────────────────────────────────────────

def migrate(
    src_path: str,
    target_theme: str,
    source_theme: str | None = None,
    output_path: str | None = None,
) -> str:
    """
    PPTX のテーマを移行する。

    Args:
        src_path:     移行元 PPTX
        target_theme: 移行先テーマ名（themes/*.json の name フィールド）
        source_theme: 移行元テーマ名（省略時は自動検出）
        output_path:  出力先（省略時は <input>_<target>.pptx）

    Returns:
        出力ファイルパス
    """
    if output_path is None:
        base = os.path.splitext(src_path)[0]
        output_path = f"{base}_{target_theme}.pptx"

    print(f"Source : {src_path}")
    print(f"Theme  : {source_theme or 'auto-detect'} → {target_theme}")
    print(f"Output : {output_path}")
    print()

    # ── テーマ読み込み ──
    dst_data     = _load_theme_data(target_theme)
    dst_tokens   = dst_data.get("tokens", {})
    dst_tpl_path = _resolve_template_path(dst_data)
    print(f"New template: {dst_tpl_path}")

    src_tokens = (
        _load_theme_data(source_theme).get("tokens", {})
        if source_theme
        else _auto_detect_source_theme(src_path)
    )

    # ── プレゼンテーション ──
    src_prs = Presentation(src_path)
    dst_prs = Presentation(dst_tpl_path)

    # ソーステーマの schemeClr マップを構築
    src_theme_map = _build_src_theme_map(src_prs)

    # スライドで実際に使われている色を収集（近似マッチに使用）
    actual_colors = _collect_slide_colors(src_prs)

    color_map = _build_color_map_fuzzy(src_tokens, dst_tokens, actual_colors)
    if color_map:
        print(f"\nColor map ({len(color_map)} replacements):")
        for old, new in sorted(color_map.items()):
            print(f"  #{old} → #{new}")
    else:
        print("\nNo color replacements (themes may use identical colors or detection failed)")
    print()

    # 新テンプレートの既存スライドをクリア
    while len(dst_prs.slides):
        sld_id = dst_prs.slides._sldIdLst[0]
        r_id   = sld_id.get(qn("r:id"))
        dst_prs.part.drop_rel(r_id)
        del dst_prs.slides._sldIdLst[0]

    # レイアウトタイプ別マップ（新テンプレート）
    dst_layout_map = _build_dst_layout_map(dst_prs)

    # ── スライドをコピー ──
    print(f"Migrating {len(src_prs.slides)} slides...")
    for i, src_slide in enumerate(src_prs.slides, 1):
        src_layout_name = src_slide.slide_layout.name
        src_layout_type = _classify_layout_type(src_layout_name)
        new_layout      = dst_layout_map.get(src_layout_type, dst_layout_map["content"])
        print(f"  Slide {i:2d}: [{src_layout_name}] → {src_layout_type} → [{new_layout.name}]")

        dst_slide = dst_prs.slides.add_slide(new_layout)

        # dst レイアウトが持つ placeholder idx セット
        dst_ph_idxs = {
            ph.placeholder_format.idx
            for ph in dst_slide.placeholders
        }

        # コンテンツスライドはタイトル (idx=0) のみプレースホルダー経由でコピー。
        # idx=10 (body) / idx=11 (breadcrumb) はテンプレート間で位置が異なるため
        # free-floating textbox として元の位置を保持する。
        if src_layout_type == "content":
            dst_ph_idxs = dst_ph_idxs & {0}

        _copy_shapes(src_slide, dst_slide, dst_ph_idxs)
        _copy_placeholder_text(src_slide, dst_slide, dst_ph_idxs)

        # 空のプレースホルダーを削除（「タイトルを入力」等のヒントテキスト表示防止）
        _remove_empty_placeholders(dst_slide)

        # schemeClr を解決して srgbClr に変換 → color_map も適用
        if src_theme_map:
            _resolve_scheme_colors(dst_slide._element, src_theme_map, color_map)

        # srgbClr を color_map で置換（schemeClr 解決済み分 + 元から srgbClr だった分）
        replaced = 0
        if color_map:
            replaced = _replace_colors(dst_slide._element, color_map)

        replaced_str = f", {replaced} color(s) replaced" if replaced else ""
        print(f"          → OK{replaced_str}")

    # ── 保存 ──
    dst_prs.save(output_path)
    print(f"\nSaved: {output_path}")
    return output_path


def main():
    parser = argparse.ArgumentParser(
        description="PPTX テーマ移行ツール — カラーとスライドマスターを一括置き換え"
    )
    parser.add_argument("input",        help="移行元 PPTX ファイルパス")
    parser.add_argument("target_theme", help="移行先テーマ名 (例: fiori, accenture)")
    parser.add_argument("--from", dest="source_theme", metavar="THEME",
                        help="移行元テーマ名（省略時は自動検出）")
    parser.add_argument("--out",  metavar="FILE",
                        help="出力ファイルパス（省略時は input_<target>.pptx）")
    args = parser.parse_args()

    migrate(
        src_path=args.input,
        target_theme=args.target_theme,
        source_theme=args.source_theme,
        output_path=args.out,
    )


if __name__ == "__main__":
    main()
