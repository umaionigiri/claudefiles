"""
Pattern X -- Step Chart (Phased Step Process)

Horizontal step chart with phase grouping headers spanning multiple steps.
Each step has a numbered title and detailed bullet content below.

Structure:
    [   Phase A   ][     Phase B     ][     Phase C     ]   <-- phase header row
    [ 1. Step ][ 2. Step ][ 3. Step ][ 4. Step ][ 5. Step ] <-- step header row
    [ bullets ][ bullets ][ bullets ][ bullets ][ bullets ] <-- detail area

Usage:
    import sys, os
    sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
    from pattern_x import add_step_chart
    from helpers import *

    phases = [
        {
            "label": "分析",
            "steps": [
                {"title": "市場分析", "subtitle": "市場ニーズ",
                 "bullets": ["市場のニーズを徹底的に調査", "関連する規制や法律を調査"]},
            ],
        },
        {
            "label": "構築",
            "steps": [
                {"title": "システム構築", "subtitle": "設備投資",
                 "bullets": ["必要な設備を選定し投資計画を立てる", "新しい設備の導入に伴う人材育成"]},
                {"title": "セキュリティ", "subtitle": "リスク管理",
                 "bullets": ["サイバーセキュリティのリスクを評価", "セキュリティ強化のための研修"]},
            ],
        },
    ]
    add_step_chart(slide, phases)
"""

import os
import sys

_SCRIPTS_DIR = os.path.dirname(os.path.abspath(__file__))
if _SCRIPTS_DIR not in sys.path:
    sys.path.insert(0, _SCRIPTS_DIR)

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.oxml.ns import qn as _qn


def add_step_chart(
    slide,
    phases,
    x=None,
    y=None,
    total_w=None,
    total_h=None,
    phase_h=0.45,
    step_header_h=0.50,
    font_name=None,
):
    """
    Draw a phased step chart on the given slide.

    Args:
        slide       : python-pptx Slide object
        phases      : list of dicts, each with:
                        'label' (str) -- phase group name (e.g. "分析")
                        'color' (RGBColor, optional) -- phase header fill color
                        'steps' (list of dicts):
                            'title'    (str) -- step title (e.g. "市場分析")
                            'subtitle' (str, optional) -- bold subtitle above bullets
                            'bullets'  (list[str]) -- detail bullet texts
        x           : left edge in inches (default: helpers.ML)
        y           : top edge in inches (default: helpers.CY)
        total_w     : total width in inches (default: helpers.CW)
        total_h     : total height in inches (default: helpers.BY - helpers.CY)
        phase_h     : height of phase header row in inches (default 0.45)
        step_header_h: height of step number+title row in inches (default 0.50)
        font_name   : font family (default: helpers.FONT)

    Layout (vertical zones from y):
        y                  -> phase header row (phase_h)
        y + phase_h        -> step header row (step_header_h)
        y + phase_h + step_header_h -> detail area (remaining height)
    """
    from helpers import (
        ML, CW, CY, BY,
        DARKEST_PURPLE, DARK_PURPLE, CORE_PURPLE, LIGHTEST_PURPLE,
        OFF_WHITE, WHITE, BLACK, TEXT_BODY, MID_GRAY, LIGHT_GRAY,
        FONT,
        accent_color,
    )

    _x = x if x is not None else ML
    _y = y if y is not None else CY
    _w = total_w if total_w is not None else CW
    _font = font_name if font_name is not None else FONT

    _by = BY - 0.10
    _h = total_h if total_h is not None else (_by - _y)

    # Count total steps across all phases
    total_steps = sum(len(ph.get("steps", [])) for ph in phases)
    if total_steps == 0:
        return

    # Column width per step
    gap = 0.08  # gap between step columns
    col_w = (_w - gap * (total_steps - 1)) / total_steps

    # Vertical zones
    phase_top = _y
    step_top = _y + phase_h
    detail_top = step_top + step_header_h
    detail_h = _h - phase_h - step_header_h

    # ── Phase header row ────────────────────────────────────────────────
    step_idx = 0
    for ph_i, phase in enumerate(phases):
        steps = phase.get("steps", [])
        n_steps = len(steps)
        if n_steps == 0:
            continue

        # Phase bar spans all its steps
        ph_x = _x + step_idx * (col_w + gap)
        ph_w = n_steps * col_w + (n_steps - 1) * gap
        ph_color = phase.get("color", DARKEST_PURPLE)

        bar = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(ph_x), Inches(phase_top),
            Inches(ph_w), Inches(phase_h),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = ph_color
        bar.line.fill.background()

        # Phase label text
        tf = bar.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = phase["label"]
        run.font.name = _font
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = WHITE
        tf._txBody.find(_qn("a:bodyPr")).set("anchor", "ctr")

        step_idx += n_steps

    # ── Step header row + detail columns ────────────────────────────────
    step_idx = 0
    for phase in phases:
        steps = phase.get("steps", [])
        for s_i, step in enumerate(steps):
            cx = _x + step_idx * (col_w + gap)
            num = step_idx + 1

            # Step header box (numbered title)
            hdr = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                Inches(cx), Inches(step_top),
                Inches(col_w), Inches(step_header_h),
            )
            hdr.fill.solid()
            hdr.fill.fore_color.rgb = LIGHTEST_PURPLE
            hdr.line.fill.background()

            tf_hdr = hdr.text_frame
            tf_hdr.word_wrap = False
            p_hdr = tf_hdr.paragraphs[0]
            p_hdr.alignment = PP_ALIGN.CENTER

            # Number run (DARK_PURPLE bold)
            r_num = p_hdr.add_run()
            r_num.text = f"{num}. "
            r_num.font.name = _font
            r_num.font.size = Pt(14)
            r_num.font.bold = True
            r_num.font.color.rgb = DARK_PURPLE

            # Title run
            r_title = p_hdr.add_run()
            r_title.text = step["title"]
            r_title.font.name = _font
            r_title.font.size = Pt(14)
            r_title.font.bold = True
            r_title.font.color.rgb = DARK_PURPLE

            tf_hdr._txBody.find(_qn("a:bodyPr")).set("anchor", "ctr")

            # ── Detail area ─────────────────────────────────────────────
            # Background panel
            panel = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                Inches(cx), Inches(detail_top),
                Inches(col_w), Inches(detail_h),
            )
            panel.fill.solid()
            panel.fill.fore_color.rgb = OFF_WHITE
            panel.line.color.rgb = LIGHT_GRAY
            panel.line.width = Pt(0.5)

            # Text content
            pad = 0.10
            tb = slide.shapes.add_textbox(
                Inches(cx + pad), Inches(detail_top + pad),
                Inches(col_w - pad * 2), Inches(detail_h - pad * 2),
            )
            tf = tb.text_frame
            tf.word_wrap = True

            line_idx = 0

            # Optional subtitle (bold header inside detail)
            subtitle = step.get("subtitle", "")
            if subtitle:
                p_sub = tf.paragraphs[0] if line_idx == 0 else tf.add_paragraph()
                run_sub = p_sub.add_run()
                run_sub.text = f"\u25cf{subtitle}"
                run_sub.font.name = _font
                run_sub.font.size = Pt(14)
                run_sub.font.bold = True
                run_sub.font.color.rgb = BLACK
                line_idx += 1

            # Bullets
            bullets = step.get("bullets", [])
            for b_text in bullets:
                p_b = tf.paragraphs[0] if line_idx == 0 else tf.add_paragraph()
                run_b = p_b.add_run()
                # Add bullet prefix if not already present
                text = b_text if b_text.startswith(("\u2022", "\u25cf", "\u30fb", "-")) else f"\u30fb{b_text}"
                run_b.text = text
                run_b.font.name = _font
                run_b.font.size = Pt(14)
                run_b.font.color.rgb = TEXT_BODY
                line_idx += 1

            step_idx += 1
