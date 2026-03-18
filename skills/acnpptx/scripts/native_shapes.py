"""
PPTX Skill — Native Shape Helpers

Replaces SVG-based decorations with python-pptx native objects where possible.
No Node.js / @resvg/resvg-js required for these functions.

Usage in generation scripts:
    sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
    from native_shapes import *

All position/size arguments are in inches unless noted otherwise.
All color arguments accept RGBColor instances (from pptx.dml.color).
"""

from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree

# ── Default colors — imported from helpers.py for theme consistency ────────────
# These are module-level references. When load_theme() changes helpers.py globals,
# functions below read the current value via helpers module at call time.
import helpers as _h

def _get_purple():
    return _h.CORE_PURPLE

def _get_deep_purple():
    return _h.DARKEST_PURPLE

_WHITE       = RGBColor(0xFF, 0xFF, 0xFF)   # never changes across themes

def _get_light_gray():
    return _h.LIGHT_GRAY


# ─────────────────────────────────────────────────────────────────────────────
# INTERNAL HELPERS
# ─────────────────────────────────────────────────────────────────────────────

def _rgb_hex(color: RGBColor) -> str:
    """Return 6-char uppercase hex string from RGBColor."""
    return f"{color[0]:02X}{color[1]:02X}{color[2]:02X}"


def _relative_luminance(color: RGBColor) -> float:
    """Calculate relative luminance per WCAG 2.0 (0.0=black, 1.0=white)."""
    def _linearize(v):
        s = v / 255.0
        return s / 12.92 if s <= 0.04045 else ((s + 0.055) / 1.055) ** 2.4
    r, g, b = _linearize(color[0]), _linearize(color[1]), _linearize(color[2])
    return 0.2126 * r + 0.7152 * g + 0.0722 * b


def _contrast_ratio(l1: float, l2: float) -> float:
    """WCAG contrast ratio between two luminances."""
    lighter = max(l1, l2)
    darker = min(l1, l2)
    return (lighter + 0.05) / (darker + 0.05)


def auto_text_color(bg_color: RGBColor) -> RGBColor:
    """
    Return white or black text color based on which has higher contrast
    against the given background color (WCAG relative luminance).
    """
    bg_lum = _relative_luminance(bg_color)
    white_contrast = _contrast_ratio(1.0, bg_lum)   # luminance of white = 1.0
    black_contrast = _contrast_ratio(bg_lum, 0.0)   # luminance of black = 0.0
    return RGBColor(0xFF, 0xFF, 0xFF) if white_contrast >= black_contrast else RGBColor(0x00, 0x00, 0x00)


def _add_preset_shape(slide, prst: str, x, y, w, h, fill_color: RGBColor,
                      line_color=None, line_width_pt=0.0, flip_h=False, flip_v=False,
                      rotation_deg=None):
    """
    Low-level: insert a preset geometry shape directly via lxml.

    Bypasses python-pptx enum mapping issues (e.g. CHEVRON producing wrong shape).
    prst must be a valid OOXML preset geometry name:
        "pentagon", "chevron", "rightArrow", "downArrow",
        "rectangle", "ellipse", "diamond", "triangle", etc.

    rotation_deg: clockwise rotation in degrees (e.g. 90 = point right for triangle).
        Converted to OOXML counterclockwise units (1/60000 degree).
    """
    spTree = slide.shapes._spTree

    # Unique ID: use max_shape_id + 1 to avoid collisions with python-pptx's own ID counter
    sp_id = str(slide.shapes._spTree.max_shape_id + 1)

    sp = etree.SubElement(spTree, qn('p:sp'))

    # nvSpPr
    nvSpPr = etree.SubElement(sp, qn('p:nvSpPr'))
    cNvPr  = etree.SubElement(nvSpPr, qn('p:cNvPr'))
    cNvPr.set('id', sp_id)
    cNvPr.set('name', f'{prst}{sp_id}')
    cNvSpPr = etree.SubElement(nvSpPr, qn('p:cNvSpPr'))
    etree.SubElement(nvSpPr, qn('p:nvPr'))

    # spPr
    spPr = etree.SubElement(sp, qn('p:spPr'))

    xfrm = etree.SubElement(spPr, qn('a:xfrm'))
    if flip_h:
        xfrm.set('flipH', '1')
    if flip_v:
        xfrm.set('flipV', '1')
    if rotation_deg is not None:
        # OOXML rotation is counterclockwise in units of 1/60000 degree.
        # Clockwise rotation_deg → negate, then mod 21600000.
        rot = int(-rotation_deg * 60000) % 21600000
        xfrm.set('rot', str(rot))
    off = etree.SubElement(xfrm, qn('a:off'))
    off.set('x', str(int(Inches(x))))
    off.set('y', str(int(Inches(y))))
    ext = etree.SubElement(xfrm, qn('a:ext'))
    ext.set('cx', str(int(Inches(w))))
    ext.set('cy', str(int(Inches(h))))

    prstGeom = etree.SubElement(spPr, qn('a:prstGeom'))
    prstGeom.set('prst', prst)
    etree.SubElement(prstGeom, qn('a:avLst'))

    # fill
    solidFill = etree.SubElement(spPr, qn('a:solidFill'))
    srgbClr   = etree.SubElement(solidFill, qn('a:srgbClr'))
    srgbClr.set('val', _rgb_hex(fill_color))

    # line
    ln = etree.SubElement(spPr, qn('a:ln'))
    if line_color and line_width_pt > 0:
        ln.set('w', str(int(Pt(line_width_pt))))
        lnFill = etree.SubElement(ln, qn('a:solidFill'))
        lnClr  = etree.SubElement(lnFill, qn('a:srgbClr'))
        lnClr.set('val', _rgb_hex(line_color))
    else:
        etree.SubElement(ln, qn('a:noFill'))

    # empty txBody (required for valid p:sp)
    txBody = etree.SubElement(sp, qn('p:txBody'))
    bodyPr = etree.SubElement(txBody, qn('a:bodyPr'))
    # Vertically center text and set tight margins to prevent overflow
    bodyPr.set('anchor', 'ctr')
    bodyPr.set('lIns', '91440')   # 0.1" left
    bodyPr.set('tIns', '45720')   # 0.05" top
    bodyPr.set('rIns', '91440')   # 0.1" right
    bodyPr.set('bIns', '45720')   # 0.05" bottom
    etree.SubElement(txBody, qn('a:lstStyle'))
    etree.SubElement(txBody, qn('a:p'))

    return sp


def _set_connector_arrowhead(connector, arrow_end=True, arrow_start=False,
                              head_type='arrow', tail_type='arrow'):
    """
    Set arrowheads on a python-pptx Connector via lxml.

    connector: result of slide.shapes.add_connector(...)
    arrow_end: put arrowhead at the END (tail in OOXML terms)
    arrow_start: put arrowhead at the START (head in OOXML terms)
    """
    ln = connector.line._ln  # a:ln element

    # Remove existing headEnd/tailEnd if any
    for tag in ('a:headEnd', 'a:tailEnd'):
        for el in ln.findall(qn(tag)):
            ln.remove(el)

    headEnd = etree.SubElement(ln, qn('a:headEnd'))
    headEnd.set('type', head_type if arrow_start else 'none')

    tailEnd = etree.SubElement(ln, qn('a:tailEnd'))
    tailEnd.set('type', tail_type if arrow_end else 'none')


# ─────────────────────────────────────────────────────────────────────────────
# PUBLIC API — ARROWS
# ─────────────────────────────────────────────────────────────────────────────

def add_arrow_right(slide, x, y, w, h, color: RGBColor = None, outline_only=False):
    """
    Add a right-pointing solid arrow shape (native OOXML rightArrow).

    Replaces: svg_arrow_right / svg_arrow_right_outline

    Args:
        x, y, w, h : position and size in inches
        color       : RGBColor (default: purple)
        outline_only: if True, no fill + colored border (outline arrow)
    Returns:
        lxml element of the inserted shape
    """
    color = color or _get_purple()
    if outline_only:
        return _add_preset_shape(slide, 'rightArrow', x, y, w, h,
                                 fill_color=_WHITE,
                                 line_color=color, line_width_pt=1.5)
    return _add_preset_shape(slide, 'rightArrow', x, y, w, h, fill_color=color)


def add_arrow_left(slide, x, y, w, h, color: RGBColor = None):
    """Left-pointing solid arrow (flipped rightArrow)."""
    color = color or _get_purple()
    return _add_preset_shape(slide, 'rightArrow', x, y, w, h,
                             fill_color=color, flip_h=True)


def add_arrow_down(slide, x, y, w, h, color: RGBColor = None):
    """
    Add a downward-pointing solid arrow shape.

    Replaces: svg_arrow_down
    """
    color = color or _get_purple()
    return _add_preset_shape(slide, 'downArrow', x, y, w, h, fill_color=color)


def add_arrow_up(slide, x, y, w, h, color: RGBColor = None):
    """Upward-pointing solid arrow."""
    color = color or _get_purple()
    return _add_preset_shape(slide, 'upArrow', x, y, w, h, fill_color=color)


def add_step_connector(slide, x, y, size=0.20, color: RGBColor = None):
    """
    Small right-pointing arrow used as step separator between process columns.

    Replaces: svg_step_connector

    Args:
        x, y  : center position in inches
        size  : arrow width and height (square, default 0.20")
    """
    color = color or _get_purple()
    return _add_preset_shape(slide, 'rightArrow',
                             x - size / 2, y - size / 2, size, size,
                             fill_color=color)


# ─────────────────────────────────────────────────────────────────────────────
# PUBLIC API — CONNECTORS (LINES WITH ARROWHEADS)
# ─────────────────────────────────────────────────────────────────────────────

def add_connector_arrow(slide, x1, y1, x2, y2,
                        color: RGBColor = None,
                        width_pt=1.5,
                        arrow_end=True,
                        arrow_start=False,
                        connector_type='straight'):
    """
    Add a native connector line with optional arrowheads.

    Replaces: svg_connector_elbow / svg_arrow_curved (as native object)

    Args:
        x1, y1     : start point in inches
        x2, y2     : end point in inches
        color      : RGBColor (default: purple)
        width_pt   : line width in points (default 1.5pt)
        arrow_end  : arrowhead at end (default True)
        arrow_start: arrowhead at start (default False)
        connector_type: 'straight' | 'elbow' | 'curved'
    Returns:
        python-pptx Connector object
    """
    from pptx.enum.shapes import MSO_CONNECTOR_TYPE

    color = color or _get_purple()

    type_map = {
        'straight': MSO_CONNECTOR_TYPE.STRAIGHT,
        'elbow'   : MSO_CONNECTOR_TYPE.ELBOW,
        'curved'  : MSO_CONNECTOR_TYPE.CURVE,
    }
    ctype = type_map.get(connector_type, MSO_CONNECTOR_TYPE.STRAIGHT)

    conn = slide.shapes.add_connector(
        ctype,
        Inches(x1), Inches(y1),
        Inches(x2), Inches(y2)
    )

    # Line style
    conn.line.width = Pt(width_pt)
    conn.line.color.rgb = color

    # Arrowheads via lxml
    _set_connector_arrowhead(conn, arrow_end=arrow_end, arrow_start=arrow_start)

    return conn


def add_divider_line(slide, x, y, w, color: RGBColor = None, width_pt=1.0):
    """
    Add a horizontal divider line (native straight connector, no arrowhead).

    Replaces: svg_divider_line (simple version without diamond accent)

    Args:
        x, y : start position in inches
        w    : line length in inches
        color: RGBColor (default: LIGHT_GRAY)
    """
    color = color or _get_light_gray()
    from pptx.enum.shapes import MSO_CONNECTOR_TYPE as _MCT
    conn = slide.shapes.add_connector(
        _MCT.STRAIGHT,
        Inches(x), Inches(y),
        Inches(x + w), Inches(y)
    )
    conn.line.width = Pt(width_pt)
    conn.line.color.rgb = color
    _set_connector_arrowhead(conn, arrow_end=False, arrow_start=False)
    return conn


# ─────────────────────────────────────────────────────────────────────────────
# PUBLIC API — CHEVRON FLOW
# ─────────────────────────────────────────────────────────────────────────────

def add_chevron_flow(slide, items, x, y, total_w, h,
                     gap=0.10,
                     fill_color: RGBColor = None,
                     text_color: RGBColor = None,
                     font_name='Meiryo UI',
                     font_size_pt=14,
                     use_pentagon_first=True,
                     shape_style='chevron'):
    """
    Add a horizontal process flow using native OOXML shapes.

    Args:
        slide            : target slide
        items            : list of strings (step labels)
        x, y             : top-left of the entire flow in inches
        total_w          : total width of the flow in inches
        h                : height of each shape in inches
        gap              : horizontal gap (inches, default 0.05)
        fill_color       : RGBColor for shape fill (default: DEEP_PURPLE)
        text_color       : RGBColor for label text (default: WHITE)
        font_name        : font family
        font_size_pt     : label font size in points (min 14; default 14)
        use_pentagon_first: if True (default), first item uses homePlate (flat left side ▷).
                           Remaining items use chevron (concave left ▷). This is the REQUIRED style.
                           homePlate = process-flow first step shape (flat left, pointed right).
                           Never set to False — all chevron flows must use homePlate as first shape.
        shape_style      : 'chevron' — pentagon(first) + chevron(rest) shapes (default, REQUIRED).
                           'box_triangle' — rectangles with right-pointing triangle separators ▷
        NOTE: 交互に向きが変わるスタイル（← → ← → 等）は非サポート。常に右向きのフロー。
              複数行への分割も禁止。全アイテムを1行に収めること。
    Returns:
        list of shape elements (or (shape_el, textbox) tuples for 'chevron' style)
    """
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn as _qn
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as _MSO

    fill_color = fill_color or _get_deep_purple()
    text_color = text_color or auto_text_color(fill_color)

    n = len(items)
    if n == 0:
        return []

    if shape_style == 'box_triangle':
        # Rectangles as step boxes with right-pointing triangle separators between them.
        # Layout: [Box]  ▷  [Box]  ▷  [Box] ...
        #
        # OOXML rotation note: rot attribute is CLOCKWISE in 1/60000 degree units.
        # The formula: rot = int(-rotation_deg * 60000) % 21600000
        #   rotation_deg=-90 → rot=5400000 = 90° CW → apex of ▲ rotates right → ▶ RIGHT ✓
        #   rotation_deg=+90 → rot=16200000 = 270° CW → apex rotates left → ◀ LEFT ✗
        # NEVER use rotation_deg=+90 — it gives LEFT-pointing ◁ in actual PowerPoint.
        #
        # Triangle proportions: unrotated cx:cy = 3:1 (wide, short box).
        # After 90° CW rotation: visual height (vertical) = cx = h, visual width (horizontal) = cy = h/3.
        sep_gap = max(gap, 0.10)   # gap on each side of the separator — minimum 0.10"

        # OOXML unrotated dimensions: wide (tri_cx=h) and short (tri_cy=h/3)
        # After 90° CW: visual height = tri_cx = h, visual width = tri_cy = h/3
        tri_cx = h          # OOXML cx (wide) → becomes visual HEIGHT after 90° CW rotation
        tri_cy = h / 3      # OOXML cy (short) → becomes visual WIDTH after 90° CW rotation
        visual_sep_w = tri_cy   # horizontal extent of separator on screen = h/3

        # Vertically center the unrotated bbox so visual height (= tri_cx = h) fills the row
        sep_y = y + (h - tri_cy) / 2   # = y + h/3

        # Layout: total_w = n*box_w + (n-1)*(2*sep_gap + visual_sep_w)
        box_w = (total_w - visual_sep_w * (n - 1) - sep_gap * 2 * (n - 1)) / n

        results = []
        cx = x
        for i, label in enumerate(items):
            # Step box (rectangle)
            box = slide.shapes.add_shape(
                _MSO.RECTANGLE,
                Inches(cx), Inches(y), Inches(box_w), Inches(h))
            box.fill.solid()
            box.fill.fore_color.rgb = fill_color
            box.line.fill.background()
            tf = box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = label
            run.font.name = font_name
            run.font.size = Pt(font_size_pt)
            run.font.bold = True
            run.font.color.rgb = text_color
            bodyPr = tf._txBody.find(_qn('a:bodyPr'))
            if bodyPr is not None:
                bodyPr.set('anchor', 'ctr')
            results.append(box)

            if i < n - 1:
                # Place triangle so its VISUAL left edge = right-of-box + sep_gap.
                # After rotation, visual_center_x = nom_x + tri_cx/2.
                # visual_left = visual_center_x - visual_sep_w/2
                #             = nom_x + tri_cx/2 - tri_cy/2
                # Set visual_left = cx + box_w + sep_gap:
                #   nom_x = (cx + box_w + sep_gap) - tri_cx/2 + tri_cy/2
                #         = (cx + box_w + sep_gap) - h/2 + h/6
                #         = (cx + box_w + sep_gap) - h/3
                sep_visual_left = cx + box_w + sep_gap
                nom_x = sep_visual_left - tri_cx / 2 + tri_cy / 2

                _add_preset_shape(slide, 'triangle',
                                  nom_x, sep_y, tri_cx, tri_cy,
                                  fill_color=fill_color,
                                  rotation_deg=-90)  # -90 → 90° CW in OOXML → right-pointing ▷

                cx = sep_visual_left + visual_sep_w + sep_gap  # left edge of next box

        return results

    # Default: 'chevron' style
    # Each chevron overlaps the previous by h*0.18 for a seamless connected look.
    overlap = h * 0.18
    item_w = (total_w - gap * (n - 1) + overlap * (n - 1)) / n

    results = []
    cx = x
    for i, label in enumerate(items):
        prst = 'homePlate' if (i == 0 and use_pentagon_first) else 'chevron'
        sp = _add_preset_shape(slide, prst, cx, y, item_w, h,
                               fill_color=fill_color)

        # Text overlay (centered textbox on top of shape)
        # Chevron body is roughly the left 80% (rest is the arrowhead);
        # adjust text box to sit in that region
        txt_x = cx + (item_w * 0.05)
        txt_w = item_w * 0.85 if i < n - 1 else item_w * 0.90
        tb = slide.shapes.add_textbox(
            Inches(txt_x), Inches(y),
            Inches(txt_w), Inches(h)
        )
        tf = tb.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = label
        run.font.name = font_name
        run.font.size = Pt(font_size_pt)
        run.font.bold = True
        run.font.color.rgb = text_color

        bodyPr = tf._txBody.find(_qn('a:bodyPr'))
        if bodyPr is not None:
            bodyPr.set('anchor', 'ctr')

        results.append((sp, tb))
        # Next chevron starts at right edge minus overlap
        cx += item_w - overlap + gap

    return results


# ─────────────────────────────────────────────────────────────────────────────
# PUBLIC API — MISC NATIVE DECORATIONS
# ─────────────────────────────────────────────────────────────────────────────

def add_accent_corner(slide, x, y, w, h, color: RGBColor = None):
    """
    Add a solid rectangle as a decorative corner accent.

    Replaces: svg_accent_corner

    Args:
        x, y, w, h: position and size in inches
        color      : RGBColor (default: PURPLE)
    """
    color = color or _get_purple()
    return _add_preset_shape(slide, 'rectangle', x, y, w, h, fill_color=color)


def add_highlight_bar(slide, x, y, w, h,
                      bg_color: RGBColor = None,
                      border_color: RGBColor = None,
                      border_pt=1.5):
    """
    Add a full-width highlight/key-message bar (native rectangle).

    Replaces: svg_highlight_bar (without the small left-arrow decoration)
    Add text separately with slide.shapes.add_textbox() on top.

    Args:
        x, y, w, h  : position and size in inches
        bg_color    : fill color (default: LP_BG — light purple tint)
        border_color: border color (default: PURPLE)
        border_pt   : border width in points
    Returns:
        lxml shape element
    """
    bg_color     = bg_color     or RGBColor(0xF3, 0xE8, 0xFF)
    border_color = border_color or _get_purple()
    return _add_preset_shape(slide, 'rectangle', x, y, w, h,
                             fill_color=bg_color,
                             line_color=border_color,
                             line_width_pt=border_pt)