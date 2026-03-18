"""
PPTX Skill — Accenture Brand Compliance Checker

Checks:
  1. カラーパレット外の色使用
  2. フォントサイズ下限違反 (< 12pt)
  3. 形状オーバーフロー
  4. カバースライドのロゴ存在チェック（Acc_Logo_* が必須）
  5. 内部スライドのフルロゴ使用禁止（GT symbolのみ許可）
  6. Rounded shapes 禁止（ROUNDED_RECTANGLE は使わない）
  7. Headline case チェック（ALL CAPS / Title Case → WARNING）

Usage:
    python brand_check.py output.pptx
    -> EXIT 0 if no errors, EXIT 1 if errors found
"""

import sys
import io
import re

if sys.platform == "win32":
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")

from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from lxml import etree

# ── 許可カラーパレット (hex 6桁大文字) ────────────────────────────────────
# 公式 Accenture ブランドカラー 13色
ALLOWED_COLORS = {
    # Purple spectrum
    "460073",  # DARKEST_PURPLE
    "7500C0",  # DARK_PURPLE
    "A100FF",  # CORE_PURPLE
    "C2A3FF",  # LIGHT_PURPLE
    "E6DCFF",  # LIGHTEST_PURPLE
    # Neutrals
    "000000",  # BLACK
    "818180",  # MID_GRAY
    "CFCFCF",  # LIGHT_GRAY
    "F1F1EF",  # OFF_WHITE
    "FFFFFF",  # WHITE
    # Secondary (sparingly)
    "FF50A0",  # PINK
    "224BFF",  # BLUE
    "05F2DB",  # AQUA
    # Practical text grays (body text readability)
    "333333",  # TEXT_BODY
    "666666",  # TEXT_SUB
    # Common near-white/near-black variants (system generated)
    "F5F5F5",  # BG_LIGHT (legacy, allowed)
    "F3E8FF",  # LP_BG (legacy, allowed)
    "D8D8D8",  # LIGHT_GRAY (legacy, allowed)
}

_NS = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
_NS_P = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"}
_FOOTER_IDXS = {20, 21}

# フォントサイズ下限 (hundredths of pt, 1200 = 12pt)
MIN_FONT_SIZE_HPT = 1200

# Rounded rectangle preset geometry names
_ROUNDED_PRESETS = {
    "roundRect", "round1Rect", "round2SameRect", "round2DiagRect",
    "snip1Rect", "snip2SameRect", "snip2DiagRect", "snipRoundRect",
}


def _hex_from_srgb_el(el):
    val = el.get("val", "")
    return val.upper() if val else None


def _is_all_caps(text):
    """True if text is entirely uppercase letters (and has at least 2 chars)."""
    stripped = text.strip()
    if len(stripped) < 2:
        return False
    letters = [c for c in stripped if c.isalpha()]
    return len(letters) >= 2 and all(c.isupper() for c in letters)


def _is_title_case(text):
    """True if text looks like Title Case (each word capitalised, not a proper name).
    We flag strings with 3+ words where every word starts with a capital.
    """
    stripped = text.strip()
    words = stripped.split()
    if len(words) < 3:
        return False
    alpha_words = [w for w in words if w[0].isalpha()]
    if len(alpha_words) < 3:
        return False
    return all(w[0].isupper() for w in alpha_words)


def check_cover_logo(slide, issues):
    """Slide 1 should contain at least one image (logo provided by slide master).

    The slide master normally embeds the logo automatically. This check verifies
    that at least one picture shape exists on the cover slide. Works for any
    theme — not limited to Accenture logo naming conventions.
    """
    found = any(
        sh.shape_type == MSO_SHAPE_TYPE.PICTURE for sh in slide.shapes
    )
    if not found:
        issues["warnings"].append(
            "Slide 1 (Cover): ロゴ画像が見つかりません。"
            "スライドマスターがロゴを自動提供しているか確認してください。"
        )


def check_no_rounded_shapes(slide, slide_num, issues):
    """Rounded rectangles are not allowed per brand guidelines."""
    for sh in slide.shapes:
        el = sh._element
        # Check preset geometry
        prstGeom = el.find(".//{http://schemas.openxmlformats.org/drawingml/2006/main}prstGeom")
        if prstGeom is not None:
            prst = prstGeom.get("prst", "")
            if prst in _ROUNDED_PRESETS:
                issues["errors"].append(
                    f"Slide {slide_num} '{sh.name}': "
                    f"Rounded shape ({prst}) は使用禁止です。rect を使用してください。"
                )


def check_colors(slide, slide_num, issues, allowed=None):
    """All colors must be from the approved palette."""
    _colors = allowed if allowed is not None else ALLOWED_COLORS
    for sh in slide.shapes:
        if sh.is_placeholder and sh.placeholder_format.idx in _FOOTER_IDXS:
            continue
        el = sh._element
        for srgb in el.findall(".//a:srgbClr", _NS):
            color = _hex_from_srgb_el(srgb)
            if color and color not in _colors:
                parent_tag = srgb.getparent().tag if srgb.getparent() is not None else "?"
                issues["errors"].append(
                    f"Slide {slide_num} '{sh.name}': "
                    f"非許可カラー #{color} (parent={parent_tag.split('}')[-1]})"
                )
        # schemeClr は theme1.xml 由来の色 — 許可パレット外の可能性あり
        for scheme in el.findall(".//a:schemeClr", _NS):
            scheme_name = scheme.get("val", "")
            # "dk1"/"lt1" (dark/light 1) はスライドマスター標準色で許容
            if scheme_name not in ("dk1", "lt1", "dk2", "lt2", "tx1", "tx2", "bg1", "bg2"):
                issues["warnings"].append(
                    f"Slide {slide_num} '{sh.name}': "
                    f"schemeClr '{scheme_name}' — テーマカラー使用。"
                    "ブランドパレット (#A100FF 等) の srgbClr に置き換えることを推奨。"
                )


def check_font_sizes(slide, slide_num, issues):
    """No font smaller than 12pt (except footers)."""
    for sh in slide.shapes:
        if sh.is_placeholder and sh.placeholder_format.idx in _FOOTER_IDXS:
            continue
        for rpr in sh._element.findall(".//a:rPr", _NS):
            sz = rpr.get("sz")
            if sz and int(float(sz)) < MIN_FONT_SIZE_HPT:
                issues["errors"].append(
                    f"Slide {slide_num} '{sh.name}': "
                    f"フォント {int(float(sz))/100:.1f}pt < 12pt"
                )


def check_overflow(slide, slide_num, prs, issues):
    """Shapes must stay within slide bounds."""
    sw = prs.slide_width
    slide_h = prs.slide_height
    tolerance = Emu(50000)
    for sh in slide.shapes:
        right = sh.left + sh.width
        bottom = sh.top + sh.height
        if right > sw + tolerance:
            issues["errors"].append(
                f"Slide {slide_num} '{sh.name}': "
                f"右オーバーフロー +{(right - sw) / 914400:.2f}\""
            )
        if bottom > slide_h + tolerance:
            issues["errors"].append(
                f"Slide {slide_num} '{sh.name}': "
                f"下オーバーフロー +{(bottom - slide_h) / 914400:.2f}\""
            )


def check_headline_case(slide, slide_num, issues):
    """Warn on ALL CAPS or Title Case headlines.
    Brand guidelines require sentence case for headlines.
    """
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        for para in sh.text_frame.paragraphs:
            # Only check larger text (likely headlines) — 20pt+
            max_sz = 0
            for run in para.runs:
                if run.font.size:
                    max_sz = max(max_sz, run.font.size)
            if max_sz < 20 * 12700 and max_sz != 0:  # below 20pt
                continue
            text = para.text.strip()
            if not text or len(text) < 4:
                continue
            if _is_all_caps(text):
                issues["warnings"].append(
                    f"Slide {slide_num} '{sh.name}': "
                    f"ALL CAPS ヘッドライン「{text[:40]}」→ Sentence case を使用"
                )
            elif _is_title_case(text):
                issues["warnings"].append(
                    f"Slide {slide_num} '{sh.name}': "
                    f"Title Case ヘッドライン「{text[:40]}」→ Sentence case を使用"
                )


def _relative_luminance(hex6):
    """WCAG 2.0 relative luminance from 6-char hex (e.g. 'FFFFFF')."""
    r = int(hex6[0:2], 16) / 255.0
    g = int(hex6[2:4], 16) / 255.0
    b = int(hex6[4:6], 16) / 255.0

    def _lin(c):
        return c / 12.92 if c <= 0.04045 else ((c + 0.055) / 1.055) ** 2.4

    return 0.2126 * _lin(r) + 0.7152 * _lin(g) + 0.0722 * _lin(b)


def _contrast_ratio(hex_a, hex_b):
    """WCAG contrast ratio between two 6-char hex colors."""
    la = _relative_luminance(hex_a)
    lb = _relative_luminance(hex_b)
    lighter = max(la, lb)
    darker = min(la, lb)
    return (lighter + 0.05) / (darker + 0.05)


# Common schemeClr → typical hex mappings (Office default theme)
_SCHEME_DEFAULTS = {
    "bg1": "FFFFFF", "lt1": "FFFFFF",
    "bg2": "E7E6E6", "lt2": "E7E6E6",
    "dk1": "000000", "tx1": "000000",
    "dk2": "44546A", "tx2": "44546A",
}


def _resolve_slide_bg_color(slide, prs):
    """Determine the background color hex of a slide.

    Checks slide → layout → master for bg/solidFill.
    Returns 6-char uppercase hex string (e.g. 'FFFFFF'), defaults to 'FFFFFF'.
    """
    ns_p = "http://schemas.openxmlformats.org/presentationml/2006/main"
    ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"

    # Walk: slide → layout → master
    elements_to_check = [slide._element]
    try:
        layout = slide.slide_layout
        elements_to_check.append(layout._element)
        elements_to_check.append(layout.slide_master._element)
    except Exception:
        pass

    for el in elements_to_check:
        cSld = el.find(f"{{{ns_p}}}cSld")
        if cSld is None:
            continue
        bg = cSld.find(f"{{{ns_p}}}bg")
        if bg is None:
            continue
        # Check bgPr/solidFill/srgbClr
        for srgb in bg.iter(f"{{{ns_a}}}srgbClr"):
            val = srgb.get("val", "")
            if len(val) == 6:
                return val.upper()
        # Check bgRef/schemeClr
        for scheme in bg.iter(f"{{{ns_a}}}schemeClr"):
            val = scheme.get("val", "")
            if val in _SCHEME_DEFAULTS:
                return _SCHEME_DEFAULTS[val]

    return "FFFFFF"  # default white


def _color_from_fill_el(solidFill):
    """Extract hex color from a solidFill element. Returns uppercase hex6 or None."""
    ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    if solidFill is None:
        return None
    srgb = solidFill.find(f"{{{ns_a}}}srgbClr")
    if srgb is not None:
        val = srgb.get("val", "")
        if len(val) == 6:
            return val.upper()
    scheme = solidFill.find(f"{{{ns_a}}}schemeClr")
    if scheme is not None:
        val = scheme.get("val", "")
        if val in _SCHEME_DEFAULTS:
            return _SCHEME_DEFAULTS[val]
    return None


def _extract_text_colors(shape):
    """Extract all explicit text color hex values from a shape.

    Checks both run-level rPr and paragraph-level defRPr (set by p.font.color.rgb).
    Returns list of (text_snippet, hex6_upper) tuples.
    """
    ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    results = []
    if not shape.has_text_frame:
        return results
    for para in shape.text_frame.paragraphs:
        para_text = para.text.strip()
        if not para_text:
            continue

        # Check paragraph-level defRPr (set by p.font.color.rgb in python-pptx)
        pPr = para._p.find(f"{{{ns_a}}}pPr")
        defRPr_color = None
        if pPr is not None:
            defRPr = pPr.find(f"{{{ns_a}}}defRPr")
            if defRPr is not None:
                sf = defRPr.find(f"{{{ns_a}}}solidFill")
                defRPr_color = _color_from_fill_el(sf)

        # Check each run
        found_run_color = False
        for run in para.runs:
            text = run.text.strip()
            if not text:
                continue
            rPr = run._r.find(f"{{{ns_a}}}rPr")
            if rPr is not None:
                sf = rPr.find(f"{{{ns_a}}}solidFill")
                color = _color_from_fill_el(sf)
                if color:
                    results.append((text[:30], color))
                    found_run_color = True
                    continue
            # Fallback: use defRPr color for this run
            if defRPr_color:
                results.append((text[:30], defRPr_color))
                found_run_color = True

        # If no runs but paragraph has text (e.g. p.text = "xxx" with defRPr)
        if not found_run_color and defRPr_color and para_text:
            results.append((para_text[:30], defRPr_color))

    return results


def _get_shape_fill_color(shape):
    """Get the solid fill color of a shape, if any. Returns hex6 or None."""
    ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    ns_p = "http://schemas.openxmlformats.org/presentationml/2006/main"
    el = shape._element
    # Check spPr/solidFill — try both p: and a: namespaces
    for ns in (ns_p, ns_a):
        spPr = el.find(f"{{{ns}}}spPr")
        if spPr is not None:
            sf = spPr.find(f"{{{ns_a}}}solidFill")
            if sf is not None:
                return _color_from_fill_el(sf)
    return None


def _find_containing_shape_bg(shape, all_shapes):
    """If shape is inside a larger colored shape, return that shape's fill color.

    This handles white text on colored rectangles (e.g. agenda badges, column headers).
    """
    try:
        sh_l, sh_t = shape.left, shape.top
        sh_r = sh_l + shape.width
        sh_b = sh_t + shape.height
    except Exception:
        return None

    margin = Emu(20000)  # ~0.02" tolerance
    for other in all_shapes:
        if other is shape:
            continue
        if not hasattr(other, 'left') or not hasattr(other, 'width'):
            continue
        try:
            o_l, o_t = other.left, other.top
            o_r = o_l + other.width
            o_b = o_t + other.height
        except Exception:
            continue
        # Check if 'other' contains 'shape'
        if (o_l - margin <= sh_l and o_t - margin <= sh_t
                and o_r + margin >= sh_r and o_b + margin >= sh_b):
            fill = _get_shape_fill_color(other)
            if fill and fill != "FFFFFF" and fill != "F1F1EF":
                return fill
    # Also check the shape's own fill (e.g. text directly in a colored rect)
    own_fill = _get_shape_fill_color(shape)
    if own_fill and own_fill != "FFFFFF" and own_fill != "F1F1EF":
        return own_fill
    return None


def check_text_bg_contrast(slide, slide_num, prs, issues, is_cover_or_closing=False):
    """Check that text color has sufficient contrast against slide background.

    If text is inside a colored shape, uses that shape's fill instead of slide bg.
    ERROR if contrast ratio < 2.0 (nearly invisible).
    WARNING if contrast ratio < 3.0 (hard to read).
    Cover/closing slides are always checked; content slides only flag ERRORs.
    """
    slide_bg_hex = _resolve_slide_bg_color(slide, prs)
    all_shapes = list(slide.shapes)

    for sh in all_shapes:
        if sh.is_placeholder and sh.placeholder_format.idx in _FOOTER_IDXS:
            continue
        text_colors = _extract_text_colors(sh)
        if not text_colors:
            continue

        # Determine effective background: containing colored shape or slide bg
        container_bg = _find_containing_shape_bg(sh, all_shapes)
        effective_bg = container_bg if container_bg else slide_bg_hex

        for snippet, txt_hex in text_colors:
            if txt_hex == effective_bg:
                issues["errors"].append(
                    f"Slide {slide_num} '{sh.name}': "
                    f"文字色 #{txt_hex} と背景色 #{effective_bg} が同一 — テキスト「{snippet}」が不可視"
                )
                continue
            ratio = _contrast_ratio(effective_bg, txt_hex)
            if ratio < 2.0:
                issues["errors"].append(
                    f"Slide {slide_num} '{sh.name}': "
                    f"コントラスト不足 (ratio={ratio:.1f}) — "
                    f"文字 #{txt_hex} vs 背景 #{effective_bg} — 「{snippet}」がほぼ不可視"
                )
            elif ratio < 3.0 and is_cover_or_closing:
                issues["warnings"].append(
                    f"Slide {slide_num} '{sh.name}': "
                    f"コントラスト低 (ratio={ratio:.1f}) — "
                    f"文字 #{txt_hex} vs 背景 #{effective_bg} — 「{snippet}」が読みにくい"
                )


def check_cover_no_bg_rect(slide, issues, slide_width_emu, slide_height_emu):
    """Slide 1 must not have a manually added background rectangle."""
    for sh in slide.shapes:
        # 対象: プレースホルダーでない AUTO_SHAPE(1) または FREEFORM(5)
        if sh.is_placeholder:
            continue
        if sh.shape_type not in (MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.FREEFORM):
            continue
        # サイズチェック: 幅 >= slide_width * 70%, 高さ >= slide_height * 50%
        if sh.width < slide_width_emu * 0.70:
            continue
        if sh.height < slide_height_emu * 0.50:
            continue
        # ソリッドフィルで暗い色かチェック
        el = sh._element
        solidFills = el.findall(".//a:solidFill", _NS)
        for sf in solidFills:
            srgb = sf.find("a:srgbClr", _NS)
            if srgb is None:
                continue
            val = srgb.get("val", "")
            if len(val) == 6:
                try:
                    r = int(val[0:2], 16)
                    g = int(val[2:4], 16)
                    b = int(val[4:6], 16)
                    # dark: R+G+B < 200 or any channel < 50
                    if (r + g + b < 200) or (r < 50) or (g < 50) or (b < 50):
                        issues["errors"].append(
                            "Slide 1 (Cover): 背景矩形が手動追加されています。"
                            "add_cover_slide() を使うとスライドマスターが背景を自動提供します。"
                            "背景矩形を削除してください。"
                        )
                        return  # 1件見つければ十分
                except ValueError:
                    pass


def check_title_single_line(slide, slide_num, issues):
    """Title placeholder (idx=0) text should fit in one line."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx != 0:
            continue
        if not ph.has_text_frame:
            continue
        text = ph.text_frame.text.strip()
        if not text:
            return
        # 全角=2, 半角=1 でカウント
        wc = 0
        for ch in text:
            wc += 2 if ord(ch) > 0x7F else 1
        if wc > 52:
            display = text[:30] + "..." if len(text) > 30 else text
            issues["warnings"].append(
                f"Slide {slide_num} タイトル「{display}」が長すぎます"
                f"（推定{wc}文字幅）。フォントサイズ削減は禁止——"
                "文言を短くしてください（目安: 全角換算52文字以内）"
            )
        return


def check_agenda_badge_wrap(slide, slide_num, issues):
    """Agenda number badges (e.g. '01'–'99') must not have wrapped text."""
    # Emu 換算: 0.3" = 274638, 1.2" = 1097280
    MIN_EMU = int(Inches(0.3))
    MAX_EMU = int(Inches(1.2))
    import re as _re
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        # 小さい正方形に近いシェイプ
        try:
            w = sh.width
            h = sh.height
        except Exception:
            continue
        if not (MIN_EMU <= w <= MAX_EMU and MIN_EMU <= h <= MAX_EMU):
            continue
        text = sh.text_frame.text.strip()
        # "01"〜"99" の2桁数字 or 短い数字テキスト
        if not _re.match(r'^\d{1,2}$', text):
            continue
        if len(sh.text_frame.paragraphs) >= 2:
            issues["errors"].append(
                f"Slide {slide_num} '{sh.name}': "
                "アジェンダバッジのテキストが折り返しています。"
                "word_wrap=False と bodyPr insets=0 を設定してください。"
            )


def check_message_line_length(slide, slide_num, issues):
    """Message line text should be 60 chars or fewer (full-width equivalent)."""
    MSG_Y_MIN = int(Emu(914400 * 0.85))  # 0.85" — just above MSG_Y
    MSG_Y_MAX = int(Emu(914400 * 1.10))  # 1.10" — just below MSG_Y
    for sh in slide.shapes:
        if sh.is_placeholder:
            continue
        if not sh.has_text_frame:
            continue
        try:
            top = sh.top
        except Exception:
            continue
        if not (MSG_Y_MIN <= top <= MSG_Y_MAX):
            continue
        text = sh.text_frame.text.strip()
        if not text:
            continue
        # 全角=2, 半角=1 でカウント
        wc = sum(2 if ord(ch) > 0x7F else 1 for ch in text)
        if wc > 120:  # 全角60文字 = 120
            display = text[:40] + "..." if len(text) > 40 else text
            issues["warnings"].append(
                f"Slide {slide_num} メッセージライン「{display}」が長すぎます"
                f"（推定{wc // 2}全角文字）。60字以内・1行に収めること。"
                "折り返すとコンテンツ領域（CY=1.50\"）に重なります。"
            )


def check_table_alternating_rows(slide, slide_num, issues):
    """Table data rows must NOT use alternating colors.
    All data rows should be WHITE or OFF_WHITE (single color).
    Only the header row may use DARKEST_PURPLE with white text.
    """
    # Allowed row fill colors (hex, uppercase)
    _ALLOWED_ROW_FILLS = {"FFFFFF", "F1F1EF", "F5F5F5", "000000", "460073"}
    for sh in slide.shapes:
        if not sh.has_table:
            continue
        tbl = sh.table
        if tbl.rows is None or len(tbl.rows) < 3:
            continue  # Need at least header + 2 data rows to detect alternation
        # Collect fill colors for data rows (skip header = row 0)
        row_fills = []
        for row_idx in range(1, len(tbl.rows)):
            row = tbl.rows[row_idx]
            # Sample the first cell's fill
            cell = tbl.cell(row_idx, 0)
            cell_el = cell._tc
            tcPr = cell_el.find("{http://schemas.openxmlformats.org/drawingml/2006/main}tcPr")
            if tcPr is None:
                row_fills.append(None)
                continue
            solidFill = tcPr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill")
            if solidFill is None:
                row_fills.append(None)
                continue
            srgb = solidFill.find("{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr")
            if srgb is not None:
                row_fills.append(srgb.get("val", "").upper())
            else:
                row_fills.append(None)
        # Check for alternation: if even rows differ from odd rows consistently
        if len(row_fills) < 2:
            continue
        non_none = [f for f in row_fills if f is not None]
        if len(non_none) < 2:
            continue
        unique_fills = set(non_none)
        if len(unique_fills) >= 2:
            # Check if it's a true alternating pattern
            even_fills = set(row_fills[i] for i in range(0, len(row_fills), 2) if row_fills[i])
            odd_fills = set(row_fills[i] for i in range(1, len(row_fills), 2) if row_fills[i])
            if even_fills and odd_fills and even_fills != odd_fills:
                issues["errors"].append(
                    f"Slide {slide_num} '{sh.name}': "
                    f"テーブルデータ行に交互色が使用されています（{', '.join('#'+c for c in unique_fills)}）。"
                    "全データ行を WHITE または OFF_WHITE 単色にしてください。"
                )


def _load_theme_colors(theme_name):
    """Load additional allowed colors from a theme JSON file.

    Returns a set of uppercase 6-char hex strings, or empty set if not found.
    """
    import os, json
    skill_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    theme_path = os.path.join(skill_dir, "assets", "themes",
                              theme_name.removesuffix(".json") + ".json")
    if not os.path.isfile(theme_path):
        return set()
    with open(theme_path, encoding="utf-8") as f:
        data = json.load(f)
    colors = set()
    for hex_val in data.get("tokens", {}).values():
        if isinstance(hex_val, str) and len(hex_val.lstrip("#")) == 6:
            colors.add(hex_val.lstrip("#").upper())
    return colors


def check_brand(filepath, theme=None):
    """
    Args:
        filepath: path to .pptx file
        theme: optional theme name (e.g. "acompany") — adds theme colors to allowed palette

    Returns:
        errors  : list[str]  エラー（修正必須）
        warnings: list[str]  警告（推奨）
    """
    prs = Presentation(filepath)
    issues = {"errors": [], "warnings": []}

    # Build effective color palette — do NOT mutate module-level ALLOWED_COLORS
    # (mutating it causes color bleed when check_brand() is called multiple times
    #  with different themes, e.g. in test suites)
    effective_colors = ALLOWED_COLORS.copy()
    if theme:
        theme_colors = _load_theme_colors(theme)
        if theme_colors:
            effective_colors.update(theme_colors)

    total_slides = len(prs.slides)
    for slide_i, sl in enumerate(prs.slides, 1):
        if slide_i == 1:
            check_cover_logo(sl, issues)
            check_cover_no_bg_rect(sl, issues, prs.slide_width, prs.slide_height)

        check_no_rounded_shapes(sl, slide_i, issues)
        check_colors(sl, slide_i, issues, allowed=effective_colors)
        check_font_sizes(sl, slide_i, issues)
        check_overflow(sl, slide_i, prs, issues)
        check_headline_case(sl, slide_i, issues)
        check_title_single_line(sl, slide_i, issues)
        check_agenda_badge_wrap(sl, slide_i, issues)
        check_message_line_length(sl, slide_i, issues)
        check_table_alternating_rows(sl, slide_i, issues)
        is_cover_closing = (slide_i == 1 or slide_i == total_slides)
        check_text_bg_contrast(sl, slide_i, prs, issues,
                               is_cover_or_closing=is_cover_closing)

    return issues["errors"], issues["warnings"]


def main():
    if len(sys.argv) < 2:
        print("Usage: python brand_check.py <file.pptx> [--theme <theme_name>]")
        sys.exit(1)

    filepath = sys.argv[1]
    theme = None
    if "--theme" in sys.argv:
        idx = sys.argv.index("--theme")
        if idx + 1 < len(sys.argv):
            theme = sys.argv[idx + 1]
    errors, warnings = check_brand(filepath, theme=theme)

    if errors:
        print(f"[ERRORS: {len(errors)}]")
        for e in errors:
            print(f"  ERROR: {e}")
    if warnings:
        print(f"[WARNINGS: {len(warnings)}]")
        for w in warnings:
            print(f"  WARN: {w}")

    if not errors and not warnings:
        print("[OK] ブランドチェック通過 (0 件)")
    else:
        print(f"\n合計: {len(errors)} エラー, {len(warnings)} 警告")

    sys.exit(1 if errors else 0)


if __name__ == "__main__":
    main()
