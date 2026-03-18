"""
PPTX Skill — Native Chart Helpers (python-pptx)

Charts use the active theme's brand colors via helpers.py.
No SVG / No Node.js required — uses python-pptx's built-in chart engine.

Usage:
    from charts import add_column_chart, add_bar_chart, add_line_chart, add_pie_chart

Chart data format (series_data):
    [
        {"name": "Series A", "values": [10, 20, 30]},
        {"name": "Series B", "values": [15, 25, 35]},
    ]
"""

from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.chart.data import ChartData
from pptx.enum.chart import (
    XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION, XL_TICK_MARK
)
from lxml import etree
import helpers as _h

# ── Chart series palette — reads from helpers.py at call time ─────────────────
# When load_theme() changes colors, _get_series_colors() returns updated values.
# The 7-color order: primary, primary_dark, dark, primary_light, secondary×3

def _get_series_colors():
    """Build chart series palette from the current theme colors."""
    return [
        _h.CORE_PURPLE,      # primary
        _h.DARKEST_PURPLE,   # primary dark
        _h.DARK_PURPLE,      # dark
        _h.LIGHT_PURPLE,     # primary light
        _h.BLUE,             # secondary
        _h.AQUA,             # secondary
        _h.PINK,             # secondary
    ]

# Neutral colors — these rarely change across themes
_BLACK      = RGBColor(0x00, 0x00, 0x00)
_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)

def _get_text_body():
    return _h.TEXT_BODY

def _get_mid_gray():
    return _h.MID_GRAY

def _get_light_gray():
    return _h.LIGHT_GRAY

def _get_off_white():
    return _h.OFF_WHITE

# Default font — follows helpers.FONT (set by set_lang())
def _get_default_font():
    return _h.FONT


# ─────────────────────────────────────────────────────────────────────────────
# INTERNAL HELPERS
# ─────────────────────────────────────────────────────────────────────────────

def _rgb_hex(color: RGBColor) -> str:
    return f"{color[0]:02X}{color[1]:02X}{color[2]:02X}"


def _apply_series_colors(chart, colors=None):
    """Apply theme palette to all series in the chart."""
    colors = colors or _get_series_colors()
    for i, series in enumerate(chart.series):
        color = colors[i % len(colors)]
        fill = series.format.fill
        fill.solid()
        fill.fore_color.rgb = color


def _style_chart_axes(chart, font_name, font_size_pt=10):
    """Style value and category axes with brand fonts/colors."""
    font_name = font_name or _get_default_font()
    try:
        ax = chart.value_axis
        ax.has_title = False
        ax.tick_labels.font.name = font_name
        ax.tick_labels.font.size = Pt(font_size_pt)
        ax.tick_labels.font.color.rgb = _get_mid_gray()
        ax.tick_marks = XL_TICK_MARK.NONE
    except Exception:
        pass

    try:
        ax = chart.category_axis
        ax.tick_labels.font.name = font_name
        ax.tick_labels.font.size = Pt(font_size_pt)
        ax.tick_labels.font.color.rgb = _get_mid_gray()
        ax.tick_marks = XL_TICK_MARK.NONE
    except Exception:
        pass


def _style_chart_base(chart, title, font_name, show_legend=True, legend_pos=None):
    """Set chart title, legend, plot area, and font styling."""
    font_name = font_name or _get_default_font()
    # Chart area background: transparent
    chart.chart_style = 2  # minimal style

    # Title
    if title:
        chart.has_title = True
        chart.chart_title.has_text_frame = True
        tf = chart.chart_title.text_frame
        tf.text = title
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = _BLACK
        tf.paragraphs[0].font.name = font_name
    else:
        chart.has_title = False

    # Legend
    if show_legend:
        chart.has_legend = True
        chart.legend.position = legend_pos or XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.name = font_name
        chart.legend.font.size = Pt(12)
        chart.legend.font.color.rgb = _get_mid_gray()
    else:
        chart.has_legend = False

    # Plot area — white background via XML
    try:
        plot_area = chart.plots[0]
        # Access underlying plot XML to remove background
    except Exception:
        pass



def _make_category_chart_data(categories, series_data):
    """Build ChartData from categories and series_data list."""
    cd = ChartData()
    cd.categories = categories
    for s in series_data:
        cd.add_series(s["name"], tuple(s["values"]))
    return cd


# ─────────────────────────────────────────────────────────────────────────────
# PUBLIC API
# ─────────────────────────────────────────────────────────────────────────────

def add_column_chart(slide, title, categories, series_data,
                     x, y, w, h,
                     font_name=None,
                     show_legend=True,
                     show_data_labels=False,
                     colors=None):
    """
    Add a vertical column chart (clustered bars).

    Args:
        slide       : target slide
        title       : chart title string (or None)
        categories  : list of category labels, e.g. ["Q1", "Q2", "Q3"]
        series_data : list of {"name": str, "values": list[float]}
        x, y, w, h  : position and size in inches
        font_name   : font for labels (default Meiryo UI)
        show_legend : show legend (default True)
        show_data_labels: show value labels on bars (default False)
        colors      : list of RGBColor for series (default: current theme palette)

    Returns:
        pptx Chart object
    """
    cd = _make_category_chart_data(categories, series_data)
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(x), Inches(y), Inches(w), Inches(h),
        cd
    )
    chart = chart_frame.chart

    _apply_series_colors(chart, colors)
    _style_chart_base(chart, title, font_name, show_legend)
    _style_chart_axes(chart, font_name)

    if show_data_labels:
        for series in chart.series:
            series.has_data_labels = True
            series.data_labels.font.size = Pt(12)
            series.data_labels.font.name = font_name
            series.data_labels.font.color.rgb = _get_text_body()

    # Remove chart border


    return chart


def add_bar_chart(slide, title, categories, series_data,
                  x, y, w, h,
                  font_name=None,
                  show_legend=True,
                  show_data_labels=False,
                  colors=None):
    """
    Add a horizontal bar chart.

    Same parameters as add_column_chart.

    Returns:
        pptx Chart object
    """
    cd = _make_category_chart_data(categories, series_data)
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(x), Inches(y), Inches(w), Inches(h),
        cd
    )
    chart = chart_frame.chart

    _apply_series_colors(chart, colors)
    _style_chart_base(chart, title, font_name, show_legend)
    _style_chart_axes(chart, font_name)

    if show_data_labels:
        for series in chart.series:
            series.has_data_labels = True
            series.data_labels.font.size = Pt(12)
            series.data_labels.font.name = font_name
            series.data_labels.font.color.rgb = _get_text_body()


    return chart


def add_line_chart(slide, title, categories, series_data,
                   x, y, w, h,
                   font_name=None,
                   show_legend=True,
                   show_markers=True,
                   colors=None):
    """
    Add a line chart.

    Args:
        show_markers: show data point markers (default True)
    Returns:
        pptx Chart object
    """
    cd = _make_category_chart_data(categories, series_data)
    chart_type = (XL_CHART_TYPE.LINE_MARKERS if show_markers
                  else XL_CHART_TYPE.LINE)
    chart_frame = slide.shapes.add_chart(
        chart_type,
        Inches(x), Inches(y), Inches(w), Inches(h),
        cd
    )
    chart = chart_frame.chart

    # For line charts, color the lines instead of fills
    colors = colors or _get_series_colors()
    for i, series in enumerate(chart.series):
        color = colors[i % len(colors)]
        series.format.line.color.rgb = color
        series.format.line.width = Pt(2.0)
        if show_markers:
            try:
                series.marker.style = 4  # circle marker
                series.marker.size = 6
                series.marker.format.fill.solid()
                series.marker.format.fill.fore_color.rgb = color
            except Exception:
                pass

    _style_chart_base(chart, title, font_name, show_legend)
    _style_chart_axes(chart, font_name)

    return chart


def add_pie_chart(slide, title, labels, values,
                  x, y, w, h,
                  font_name=None,
                  show_legend=True,
                  show_data_labels=True,
                  colors=None):
    """
    Add a pie chart.

    Args:
        labels : list of category labels
        values : list of numeric values (same length as labels)
        show_data_labels: show percentages on slices (default True)
    Returns:
        pptx Chart object
    """
    cd = ChartData()
    cd.categories = labels
    cd.add_series("", tuple(values))

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE,
        Inches(x), Inches(y), Inches(w), Inches(h),
        cd
    )
    chart = chart_frame.chart

    # Color individual pie slices
    colors = colors or _get_series_colors()
    try:
        series = chart.series[0]
        for i, point in enumerate(series.points):
            color = colors[i % len(colors)]
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = color
    except Exception:
        _apply_series_colors(chart, colors)

    _style_chart_base(chart, title, font_name, show_legend,
                      legend_pos=XL_LEGEND_POSITION.RIGHT)

    if show_data_labels:
        try:
            series = chart.series[0]
            series.has_data_labels = True
            dl = series.data_labels
            dl.font.size = Pt(12)
            dl.font.name = font_name
            dl.font.color.rgb = _WHITE
            dl.position = XL_LABEL_POSITION.INSIDE_END
            dl.show_percentage = True
            dl.show_value = False
            dl.show_category_name = False
        except Exception:
            pass


    return chart


def add_stacked_column_chart(slide, title, categories, series_data,
                              x, y, w, h,
                              font_name=None,
                              show_legend=True,
                              percentage=False,
                              colors=None):
    """
    Add a stacked (or 100% stacked) column chart.

    Args:
        percentage: if True, use 100% stacked variant (default False)
    Returns:
        pptx Chart object
    """
    chart_type = (XL_CHART_TYPE.COLUMN_STACKED_100 if percentage
                  else XL_CHART_TYPE.COLUMN_STACKED)
    cd = _make_category_chart_data(categories, series_data)
    chart_frame = slide.shapes.add_chart(
        chart_type,
        Inches(x), Inches(y), Inches(w), Inches(h),
        cd
    )
    chart = chart_frame.chart

    _apply_series_colors(chart, colors)
    _style_chart_base(chart, title, font_name, show_legend)
    _style_chart_axes(chart, font_name)

    return chart


def add_area_chart(slide, title, categories, series_data,
                   x, y, w, h,
                   font_name=None,
                   show_legend=True,
                   colors=None):
    """
    Add a filled area chart.

    Returns:
        pptx Chart object
    """
    cd = _make_category_chart_data(categories, series_data)
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.AREA,
        Inches(x), Inches(y), Inches(w), Inches(h),
        cd
    )
    chart = chart_frame.chart

    _apply_series_colors(chart, colors)
    _style_chart_base(chart, title, font_name, show_legend)
    _style_chart_axes(chart, font_name)

    return chart
