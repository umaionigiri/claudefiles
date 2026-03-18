"""
Pattern V — Numbered Card Grid
Flexible N×M grid with numbered circle badges and optional emphasis borders.

Reference design: 2 rows × 4 columns (8 cards), each card has:
  - Numbered circle badge (DARK_PURPLE, top-left)
  - Bold title in header area
  - Bullet body text below
  - Optional PINK emphasis border for highlighted cards

Usage:
    import sys, os
    sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
    from pattern_v import add_numbered_card_grid
    from helpers import *

    cards = [
        {"title": "DXに特化した\nリーダーシップの確保",
         "body": "• 貴社のDXの旗振り役を務めつつ、技術スタッフの強力な巻き込みを果たすリーダーシップを確保"},
        ...
    ]
    add_numbered_card_grid(slide, cards, n_cols=4, highlight_indices=[3, 7])
"""

import os
import sys

# Allow importing helpers from the same scripts/ directory
_SCRIPTS_DIR = os.path.dirname(os.path.abspath(__file__))
if _SCRIPTS_DIR not in sys.path:
    sys.path.insert(0, _SCRIPTS_DIR)

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.oxml.ns import qn as _qn


def add_numbered_card_grid(
    slide,
    cards,
    n_cols=4,
    x=None,
    y=None,
    total_w=None,
    total_h=None,
    gap_h=0.20,
    gap_v=0.20,
    highlight_indices=None,
    highlight_color=None,
    font_name=None,
):
    """
    Draw a numbered card grid (N×M) on the given slide.

    Args:
        slide            : python-pptx Slide object
        cards            : list of dicts, each with:
                             'title' (str)  — bold heading text (use \\n for line breaks)
                             'body'  (str)  — body text (use \\n to separate bullet lines)
                             'number' (int) — optional badge number (defaults to 1-based index)
        n_cols           : number of columns (default 4)
        x                : left edge in inches (default: helpers.ML = 0.42)
        y                : top edge in inches (default: helpers.CY = 1.50)
        total_w          : total grid width in inches (default: helpers.CW = 12.50)
        total_h          : total grid height (default: fills to helpers.BY with bottom margin)
        gap_h            : horizontal gap between cards in inches (default 0.20)
        gap_v            : vertical gap between rows in inches (default 0.20)
        highlight_indices: list of 0-based card indices to draw with emphasis border
        highlight_color  : RGBColor for emphasis border (default: PINK #FF50A0)
        font_name        : font (default: helpers.FONT = "Meiryo UI")

    Layout per card:
        ┌──────────────────────────────────────────┐  ← card background (OFF_WHITE)
        │ ● badge │ TITLE (bold DARK_PURPLE 14pt)   │  ← header band (LIGHTEST_PURPLE)
        ├──────────────────────────────────────────┤
        │ • Bullet point text (TEXT_BODY 14pt)      │
        │ • More detail                             │
        └──────────────────────────────────────────┘
          ↑ optional PINK border for highlighted cards

    The badge is an ellipse (circle) with DARK_PURPLE fill and WHITE bold number text.
    """
    # ── Import helpers constants ──────────────────────────────────────────
    from helpers import (
        ML, CW, CY, BY,
        DARK_PURPLE, LIGHTEST_PURPLE, OFF_WHITE, WHITE, TEXT_BODY, PINK,
        FONT,
    )

    _x      = x        if x        is not None else ML
    _y      = y        if y        is not None else round(4.1 / 2.54, 3)  # 4.1cm from slide top
    _total_w = total_w if total_w  is not None else CW
    _hl_col  = highlight_color if highlight_color is not None else PINK
    _font    = font_name if font_name is not None else FONT
    _hl_idx  = set(highlight_indices) if highlight_indices else set()

    n_cards = len(cards)
    n_rows  = -(-n_cards // n_cols)  # ceil division

    # Derive card dimensions
    card_w = (_total_w - gap_h * (n_cols - 1)) / n_cols

    _by = BY - 0.10   # bottom margin before footer
    if total_h is not None:
        _grid_h = total_h
    else:
        _grid_h = _by - _y
    card_h = (_grid_h - gap_v * (n_rows - 1)) / n_rows

    # Header band height (badge + title zone)
    HEADER_H  = min(0.72, card_h * 0.30)
    BADGE_D   = 0.32   # badge diameter (circle)
    BADGE_PAD = 0.10   # padding from card edge to badge
    TITLE_PAD = BADGE_PAD + BADGE_D + 0.08  # x offset where title text begins

    for idx, card in enumerate(cards):
        col = idx % n_cols
        row = idx // n_cols

        cx = _x + col * (card_w + gap_h)
        cy = _y + row * (card_h + gap_v)

        num = card.get("number", idx + 1)
        title_text = card.get("title", "")
        body_text  = card.get("body",  "")
        is_hl      = idx in _hl_idx

        # ── Card background ──────────────────────────────────────────────
        bg = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(cx), Inches(cy), Inches(card_w), Inches(card_h))
        bg.fill.solid()
        bg.fill.fore_color.rgb = OFF_WHITE
        if is_hl:
            bg.line.color.rgb = _hl_col
            bg.line.width     = Pt(2.0)
        else:
            bg.line.fill.background()

        # ── Header band (LIGHTEST_PURPLE) ────────────────────────────────
        hdr_band = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(cx), Inches(cy), Inches(card_w), Inches(HEADER_H))
        hdr_band.fill.solid()
        hdr_band.fill.fore_color.rgb = LIGHTEST_PURPLE
        hdr_band.line.fill.background()

        # ── Badge (circle = OVAL) ────────────────────────────────────────
        badge_x = cx + BADGE_PAD
        badge_y = cy + (HEADER_H - BADGE_D) / 2
        badge = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            Inches(badge_x), Inches(badge_y), Inches(BADGE_D), Inches(BADGE_D))
        badge.fill.solid()
        badge.fill.fore_color.rgb = DARK_PURPLE
        badge.line.fill.background()
        tf_badge = badge.text_frame
        tf_badge.word_wrap = False
        p_badge = tf_badge.paragraphs[0]
        p_badge.text      = str(num)
        p_badge.font.size = Pt(12)
        p_badge.font.bold = True
        p_badge.font.color.rgb = WHITE
        p_badge.font.name = _font
        p_badge.alignment = PP_ALIGN.CENTER
        # vertical centre
        tf_badge._txBody.find(_qn("a:bodyPr")).set("anchor", "ctr")

        # ── Title text (in header band, right of badge) ──────────────────
        title_x = cx + TITLE_PAD
        title_w = card_w - TITLE_PAD - 0.08
        tb_title = slide.shapes.add_textbox(
            Inches(title_x), Inches(cy + 0.06),
            Inches(title_w), Inches(HEADER_H - 0.08))
        tf_t = tb_title.text_frame
        tf_t.word_wrap = True
        for li, line in enumerate(title_text.split("\n")):
            p = tf_t.paragraphs[0] if li == 0 else tf_t.add_paragraph()
            p.text            = line
            p.font.size       = Pt(14)
            p.font.bold       = True
            p.font.color.rgb  = DARK_PURPLE
            p.font.name       = _font
        tf_t._txBody.find(_qn("a:bodyPr")).set("anchor", "ctr")

        # ── Body text ────────────────────────────────────────────────────
        body_y  = cy + HEADER_H + 0.08
        body_h  = card_h - HEADER_H - 0.12
        tb_body = slide.shapes.add_textbox(
            Inches(cx + 0.10), Inches(body_y),
            Inches(card_w - 0.20), Inches(body_h))
        tf_b = tb_body.text_frame
        tf_b.word_wrap = True
        for li, line in enumerate(body_text.split("\n")):
            p = tf_b.paragraphs[0] if li == 0 else tf_b.add_paragraph()
            p.text           = line
            p.font.size      = Pt(14)
            p.font.color.rgb = TEXT_BODY
            p.font.name      = _font
