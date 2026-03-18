"""
PPTX Skill — Shared helpers, constants, and color palette.

Import this module in every slide-generation script:
    import sys, os
    sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
    from helpers import *

Language mode:
    LANG = "ja"  → Meiryo UI + slide-master.pptx  (Japanese content default)
    LANG = "en"  → Graphik + slide-master.pptx  (English content; both languages use the same template)
    Set LANG at the top of your generation script via set_lang() before importing helpers.
    Note: call set_lang() BEFORE load_theme() — set_lang resets TEMPLATE_PATH, which load_theme overwrites.

TEMPLATE_PATH and SKILL_DIR are auto-detected from this file's location,
so the generation script works from any directory.
"""

import os
import copy
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── AUTO-DETECTED PATHS (do not override) ────────────────────
_SCRIPTS_DIR = os.path.dirname(os.path.abspath(__file__))   # .../pptx/scripts/
SKILL_DIR    = os.path.dirname(_SCRIPTS_DIR)                 # .../pptx/
# ──────────────────────────────────────────────────────────────

# ── LANGUAGE / FONT MODE — SET THIS IN YOUR GENERATION SCRIPT ─
# LANG = "ja"  → Japanese: Meiryo UI + slide-master.pptx
# LANG = "en"  → English:  Graphik + slide-master.pptx (same template)
LANG = "ja"   # default; override in your script if generating English content
# ──────────────────────────────────────────────────────────────

# ── FONT CONSTANTS ────────────────────────────────────────────
FONT_JA = "Meiryo UI"   # Japanese content
FONT_EN = "Graphik"     # English content (Arial as fallback)
FONT    = FONT_JA       # active font — set via LANG below

# ── TEMPLATE PATHS ────────────────────────────────────────────
TEMPLATE_JA   = os.path.join(SKILL_DIR, "assets", "slide-master.pptx")
TEMPLATE_PATH = TEMPLATE_JA  # only Japanese template — English template removed


def set_lang(lang: str):
    """Switch font and template based on language.

    Because Python's 'from helpers import *' copies values at import time,
    this function modifies helpers module globals but NOT the caller's locals.

    Recommended usage in generation scripts:
        import helpers
        helpers.set_lang("en")        # updates helpers.FONT, helpers.TEMPLATE_PATH
        from helpers import *         # import AFTER calling set_lang
        # Or: directly assign after import
        # FONT = FONT_EN; TEMPLATE_PATH = TEMPLATE_EN

    Args:
        lang: "ja" (Meiryo UI + slide-master.pptx) or
              "en" (Graphik font + slide-master.pptx — English template removed, using Japanese master)
    """
    global LANG, FONT, TEMPLATE_PATH
    LANG = lang
    if lang == "en":
        FONT = FONT_EN
        TEMPLATE_PATH = TEMPLATE_JA  # English template removed; using Japanese master
    else:
        FONT = FONT_JA
        TEMPLATE_PATH = TEMPLATE_JA


# Logo and GT symbol assets removed — placement is handled exclusively
# by the slide master layout. Never call add_logo() or add_gt_symbol().

# ── FOOTER ────────────────────────────────────────────────────
FOOTER_TEXT = "Copyright © 2026 Accenture. All rights reserved."
# ──────────────────────────────────────────────────────────────

# ── MANDATORY CONSTANTS — DO NOT CHANGE ──────────────────────
ML  = 0.42    # margin left (inches)
MR  = 0.42    # margin right
CW  = 12.50   # content width = 13.333 - ML - MR
CY  = 1.50    # content area Y start (with message line)
BY  = 6.85    # content area Y bottom
AH  = 5.35    # available height = BY - CY
# Slide dimensions: 13.333" x 7.500" (standard 16:9)
# Layout 2 placeholder positions (y, h in inches):
#   idx=11 breadcrumb : y=0.08, h=0.27
#   idx=0  title      : y=0.42, h=0.48  (bottom=0.90)
#   Message line      : y=0.95, h=0.42  (bottom=1.37)
#   Content area      : y=1.50 → BY=6.85
# ──────────────────────────────────────────────────────────────

# ── COLOR PALETTE — Official Accenture Brand ──────────────────
#
# Purple spectrum (dark → light)
DARKEST_PURPLE  = RGBColor(0x46, 0x00, 0x73)  # #460073 — cover bg, section dividers
DARK_PURPLE     = RGBColor(0x75, 0x00, 0xC0)  # #7500C0 — emphasis, badges
CORE_PURPLE     = RGBColor(0xA1, 0x00, 0xFF)  # #A100FF — primary accent (logos, CTAs)
LIGHT_PURPLE    = RGBColor(0xC2, 0xA3, 0xFF)  # #C2A3FF — cover subtitles, tinted text
LIGHTEST_PURPLE = RGBColor(0xE6, 0xDC, 0xFF)  # #E6DCFF — subtle tinted backgrounds
#
# Neutrals
BLACK     = RGBColor(0x00, 0x00, 0x00)  # #000000 — titles
MID_GRAY  = RGBColor(0x81, 0x81, 0x80)  # #818180 — secondary text
LIGHT_GRAY = RGBColor(0xCF, 0xCF, 0xCF) # #CFCFCF — borders, dividers
OFF_WHITE = RGBColor(0xF1, 0xF1, 0xEF)  # #F1F1EF — card backgrounds, panels
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)  # #FFFFFF — text on dark, bg
#
# Secondary — use sparingly (<5% of slide area)
PINK = RGBColor(0xFF, 0x50, 0xA0)  # #FF50A0
BLUE = RGBColor(0x22, 0x4B, 0xFF)  # #224BFF
AQUA = RGBColor(0x05, 0xF2, 0xDB)  # #05F2DB
#
# Practical text grays (not in brand palette but needed for body text readability)
TEXT_BODY = RGBColor(0x33, 0x33, 0x33)  # #333333 — body text on white
TEXT_SUB  = RGBColor(0x66, 0x66, 0x66)  # #666666 — captions, footnotes
#
# Aliases for common use
PURPLE      = CORE_PURPLE      # kept for backward compatibility
DEEP_PURPLE = DARKEST_PURPLE   # kept for backward compatibility
BG_LIGHT    = OFF_WHITE        # kept for backward compatibility
# ──────────────────────────────────────────────────────────────

# ── MESSAGE LINE CONSTANTS ────────────────────────────────────
# All content slides (except cover / section / agenda) need a message line.
# Placed directly below the title placeholder (bottom=0.90").
MSG_Y  = 0.95   # message line Y position (inches, just below title)
MSG_H  = 0.45   # message line height (18pt text, bottom=1.40")
CY2    = CY     # alias: content start with message line = CY = 1.50"
# ──────────────────────────────────────────────────────────────

_slide_num = [0]


def clear_placeholders(slide):
    """Remove unused body placeholder (idx=10) from Layout 2 slide XML.

    Layout 2 has three content placeholders:
      idx=0  title      (y=0.42") — used by add_title()
      idx=10 body       (y=0.87") — NOT used; must be removed or PowerPoint
                                    shows its hint text in edit mode
      idx=11 breadcrumb (y=0.08") — used by add_breadcrumb()

    Footer/slidenum (idx=20, 21) are added later by set_footer() and are
    not present on a fresh slide.
    """
    to_remove = []
    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        if idx not in (0, 11, 20, 21):   # remove any unused body/content ph
            to_remove.append(ph._element)
    for el in to_remove:
        el.getparent().remove(el)


def add_title(slide, text, font_name=None, size_pt=28):
    """Set slide title using Layout 2's title placeholder (idx=0, y=0.42").

    The placeholder position is defined by the slide master — use this
    function instead of add_textbox() for the title.

    Args:
        slide    : target slide (Layout 2)
        text     : title text (sentence case)
        font_name: override font (defaults to module FONT)
        size_pt  : font size in points (default 28)
    """
    _font = font_name or FONT
    ph = slide.placeholders[0]
    tf = ph.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = _font
    run.font.size = Pt(size_pt)
    run.font.bold = True
    run.font.color.rgb = BLACK


def add_breadcrumb(slide, text, font_name=None):
    """Set breadcrumb text using Layout 2's small top placeholder (idx=11, y=0.08").

    Format: "Section > Topic". 12pt MID_GRAY. Pass empty string to suppress.

    Args:
        slide    : target slide (Layout 2)
        text     : breadcrumb text, e.g. "効果と実績 > KPI"
        font_name: override font (defaults to module FONT)
    """
    _font = font_name or FONT
    ph = None
    for _ph in slide.placeholders:
        if _ph.placeholder_format.idx == 11:
            ph = _ph
            break
    if ph is None:
        # Fallback: textbox at master position
        tb = slide.shapes.add_textbox(Inches(ML), Inches(0.08), Inches(CW), Inches(0.27))
        p = tb.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(12)
        p.font.color.rgb = MID_GRAY
        p.font.name = _font
        return
    tf = ph.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = _font
    run.font.size = Pt(12)
    run.font.color.rgb = MID_GRAY


def add_message_line(slide, text, font_name=None):
    """Add a message line (メッセージライン) below the slide title.

    Every content slide (except cover, section divider, agenda) should have
    one. Write 1–2 sentences summarising the slide's takeaway in 常体
    (である/だ style).

    Args:
        slide    : target slide
        text     : message text in 常体
        font_name: override font (defaults to module FONT)
    """
    _font = font_name or FONT
    tb = slide.shapes.add_textbox(
        Inches(ML), Inches(MSG_Y), Inches(CW), Inches(MSG_H))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = DARK_PURPLE
    p.font.name = _font


def set_footer(slide):
    """Clone footer/slidenum from layout and set text.

    python-pptx does NOT auto-inherit FOOTER (idx=20) and SLIDE_NUMBER
    (idx=21) placeholders from the slide layout. This function clones
    them with deepcopy and sets the configured FOOTER_TEXT.

    Call this after creating each content slide (Layout 2).
    """
    _slide_num[0] += 1
    existing = {ph.placeholder_format.idx for ph in slide.placeholders}
    if 20 not in existing or 21 not in existing:
        layout = slide.slide_layout
        for layout_ph in layout.placeholders:
            idx = layout_ph.placeholder_format.idx
            if idx in (20, 21) and idx not in existing:
                clone = copy.deepcopy(layout_ph._element)
                slide.shapes._spTree.append(clone)
    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 20:
            ph.text = FOOTER_TEXT
            for p in ph.text_frame.paragraphs:
                p.font.size = Pt(8)
                p.font.color.rgb = TEXT_SUB
                p.font.name = FONT
        elif idx == 21:
            ph.text = str(_slide_num[0])
            for p in ph.text_frame.paragraphs:
                p.font.size = Pt(8)
                p.font.color.rgb = TEXT_SUB
                p.font.name = FONT
                p.alignment = PP_ALIGN.RIGHT


def reset_slide_counter():
    """Reset the slide number counter (useful when generating multiple decks)."""
    _slide_num[0] = 0


def load_theme(theme_name_or_path: str | None = None):
    """カラーテーマを JSON から読み込んでこのモジュールの色定数を上書きする。

    テーマ名を直接指定して非対話的に適用する（推奨）。

    Args:
        theme_name_or_path: テーマ名 (e.g. "accenture") または JSON ファイルパス。
                            None の場合は何もしない（デフォルト Accenture のまま）。

    Usage (generation script の先頭):
        # テーマ名で直接指定する場合（推奨）
        import helpers as _h
        _h.load_theme("accenture")   # assets/themes/accenture.json を適用
        from helpers import *

    Note: テーマは AskUserQuestion で事前に確認し、load_theme() にハードコードすること。
          theme_selector.select_theme() は Claude Code 環境では使用禁止（Tkinter は Claude に結果を返せない）。
    """
    import json as _json
    import os as _os

    if theme_name_or_path is None:
        return

    # パス解決
    path = theme_name_or_path
    if not _os.path.isfile(path):
        # 名前として解釈 → assets/themes/<name>.json を探す
        candidate = _os.path.join(SKILL_DIR, "assets", "themes",
                                  theme_name_or_path.removesuffix(".json") + ".json")
        if _os.path.isfile(candidate):
            path = candidate
        else:
            import sys as _sys
            print(
                f"WARNING: load_theme: theme '{theme_name_or_path}' not found in "
                f"{_os.path.join(SKILL_DIR, 'assets', 'themes')} — using default Accenture colors",
                file=_sys.stderr,
            )
            return

    with open(path, encoding="utf-8") as _f:
        data = _json.load(_f)

    tokens = data.get("tokens", {})

    # トークン → モジュール変数マッピング
    _MAP = {
        "primary":       ["CORE_PURPLE", "PURPLE"],
        "primary_light": ["LIGHTEST_PURPLE"],
        "primary_dark":  ["DARKEST_PURPLE", "DEEP_PURPLE"],
        "surface":       ["OFF_WHITE", "BG_LIGHT"],
        "text_heading":  ["BLACK"],
        "text_body":     ["TEXT_BODY"],
        "text_muted":    ["MID_GRAY"],
        "border":        ["LIGHT_GRAY"],
    }
    import sys as _sys
    _mod = _sys.modules[__name__]

    for token, var_names in _MAP.items():
        _hex = tokens.get(token)
        if not _hex:
            continue
        _h = _hex.lstrip("#")
        _rgb = RGBColor(int(_h[0:2], 16), int(_h[2:4], 16), int(_h[4:6], 16))
        for _vn in var_names:
            if hasattr(_mod, _vn):
                setattr(_mod, _vn, _rgb)

    # DARK_PURPLE = primary と primary_dark の中間色（派生）
    if "primary" in tokens and "primary_dark" in tokens:
        def _mid(h1, h2):
            a, b = h1.lstrip("#"), h2.lstrip("#")
            return RGBColor(
                (int(a[0:2], 16) + int(b[0:2], 16)) // 2,
                (int(a[2:4], 16) + int(b[2:4], 16)) // 2,
                (int(a[4:6], 16) + int(b[4:6], 16)) // 2,
            )
        _mod.DARK_PURPLE = _mid(tokens["primary"], tokens["primary_dark"])

    # LIGHT_PURPLE = primary と primary_light の中間色（派生）
    if "primary" in tokens and "primary_light" in tokens:
        def _mid2(h1, h2):
            a, b = h1.lstrip("#"), h2.lstrip("#")
            return RGBColor(
                (int(a[0:2], 16) + int(b[0:2], 16)) // 2,
                (int(a[2:4], 16) + int(b[2:4], 16)) // 2,
                (int(a[4:6], 16) + int(b[4:6], 16)) // 2,
            )
        _mod.LIGHT_PURPLE = _mid2(tokens["primary"], tokens["primary_light"])

    # フッターテキスト更新（theme JSON に "footer_text" フィールドがある場合）
    footer_val = data.get("footer_text")
    if footer_val:
        _mod.FOOTER_TEXT = footer_val

    # テンプレートパス更新（theme JSON に "template" フィールドがある場合）
    template_val = data.get("template")
    if template_val:
        if _os.path.isabs(template_val):
            tpl_path = template_val
        else:
            tpl_path = _os.path.join(SKILL_DIR, template_val)
        if _os.path.isfile(tpl_path):
            _mod.TEMPLATE_PATH = tpl_path


def add_image_fit(slide, img_path, x, y, max_w, max_h):
    """Add image preserving aspect ratio within max_w x max_h bounding box (inches)."""
    from PIL import Image as PILImage
    if not os.path.isfile(img_path):
        raise FileNotFoundError(f"Image not found: {img_path}")
    pil = PILImage.open(img_path)
    iw, ih = pil.size
    ar = iw / ih
    if (max_w / max_h) > ar:   # box wider → constrain by height
        h = max_h; w = h * ar
    else:                       # box taller → constrain by width
        w = max_w; h = w / ar
    slide.shapes.add_picture(img_path, Inches(x), Inches(y), Inches(w), Inches(h))
    return w, h


# add_logo() and add_gt_symbol() removed.
# Logo and GT symbol are provided by the slide master — never add them manually.


# ── ACCENT COLOR HELPER ──────────────────────────────────────────────
# Dynamically builds a trio from the current theme colors (dark → primary → light).
# Always reflects the active theme set by load_theme(), not hardcoded values.


def accent_color(index):
    """Get accent color by index, cycling through the theme's color trio.

    The trio is built dynamically from the current module-level color
    variables (DARKEST_PURPLE, CORE_PURPLE, LIGHT_PURPLE), so it
    automatically follows load_theme() changes.

    Usage in charts/cards with multiple categories:
        for i, category in enumerate(categories):
            color = accent_color(i)          # RGBColor
            hex_c = accent_color_hex(i)      # "#460073" etc.
    """
    import sys as _sys
    _mod = _sys.modules[__name__]
    palette = [_mod.DARKEST_PURPLE, _mod.CORE_PURPLE, _mod.LIGHT_PURPLE]
    return palette[index % len(palette)]


def accent_color_hex(index):
    """Get accent color as hex string, cycling through the theme's color trio.

    Returns a "#RRGGBB" string for use in SVG fills or string formatting.
    """
    c = accent_color(index)
    return f"#{c}"


# ── MODERN VISUAL HELPERS ─────────────────────────────────────────────
# Functions that leverage the slide master for consistent branding.
# Closing slides and section dividers reuse existing layouts so that
# logo, GT symbol, and background are provided automatically.


def make_dark_divider(slide, title, subtitle=None, font_name=None):
    """Convert a slide into a dark section divider.

    Creates a full-slide DARKEST_PURPLE background with centered white title.
    This covers the slide master's elements — use sparingly (1-2 per deck max).

    Args:
        slide: target slide (any layout)
        title: section title text
        subtitle: optional smaller subtitle below
        font_name: override font (defaults to module FONT)
    """
    from pptx.enum.shapes import MSO_SHAPE
    _font = font_name or FONT
    # Full-slide dark rectangle
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(13.333), Inches(7.50),
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARKEST_PURPLE
    bg.line.fill.background()
    # Send to back
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)
    # Centered title
    tb = slide.shapes.add_textbox(Inches(ML), Inches(2.50), Inches(CW), Inches(2.00))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.size = Pt(36)
    run.font.color.rgb = WHITE
    run.font.bold = True
    run.font.name = _font
    if subtitle:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(12)
        run2 = p2.add_run()
        run2.text = subtitle
        run2.font.size = Pt(18)
        run2.font.color.rgb = LIGHT_GRAY
        run2.font.name = _font


def add_title_accent_line(slide, color=None, y=None, width=None):
    """Add a thin purple accent line below the title/message area.

    Default position: y = CY - 0.03 (just above content area).
    Call on content slides after setting the title for a modern look.

    Args:
        slide: target slide
        color: line color (defaults to CORE_PURPLE)
        y: vertical position in inches (defaults to CY - 0.03)
        width: line width in inches (defaults to CW)
    """
    from pptx.enum.shapes import MSO_SHAPE
    if color is None:
        color = CORE_PURPLE
    if y is None:
        y = CY - 0.03
    if width is None:
        width = CW
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(ML), Inches(y), Inches(width), Inches(0.03),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()


def make_closing_slide(prs, title="Thank You", subtitle=None, font_name=None,
                       layout_idx=0, text_color=None):
    """Add a closing slide using the slide master's cover layout.

    Reuses the cover layout (layout_idx=0 by default) so that background,
    logo, and GT symbol are provided automatically by the slide master.
    All placeholders are filled with the closing text — no background
    rectangle is added manually.

    The text color defaults to WHITE (for dark-background themes). Pass
    text_color=BLACK for white-background themes, or check the theme JSON's
    cover_text_color field.

    Usage:
        make_closing_slide(prs)                        # default "Thank You"
        make_closing_slide(prs, "ありがとうございました")  # Japanese
        make_closing_slide(prs, "Gracias", "¿Preguntas?")
        make_closing_slide(prs, text_color=BLACK)      # white bg theme

    Returns the created slide.
    """
    _font = font_name or FONT
    _color = text_color or WHITE
    _sub_color = LIGHT_GRAY if _color == WHITE else MID_GRAY
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    # Fill all placeholders — master provides background, logo, GT
    to_remove = []
    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        ph.text_frame.clear()
        p = ph.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        if idx == 0:  # Main title placeholder
            run = p.add_run()
            run.text = title
            run.font.size = Pt(44)
            run.font.bold = True
            run.font.color.rgb = _color
            run.font.name = _font
        elif idx == 1 and subtitle:  # Subtitle placeholder
            run = p.add_run()
            run.text = subtitle
            run.font.size = Pt(20)
            run.font.color.rgb = _sub_color
            run.font.name = _font
        else:
            # Collect unused placeholders for removal (p.text="" still shows layout hints)
            to_remove.append(ph._element)

    # Remove unused placeholders from slide XML to prevent hint text ghosts
    spTree = slide.shapes._spTree
    for sp in to_remove:
        if sp.getparent() is not None:
            spTree.remove(sp)

    return slide


def strip_sections(prs):
    """Remove all PowerPoint sections from the presentation.

    The template's presentation.xml may contain a <p14:sectionLst> element
    that creates section headers in PowerPoint's slide sorter view.
    Call this right before prs.save() to produce a clean deck.
    """
    from lxml.etree import QName
    prs_elem = prs.part._element
    ns14 = "http://schemas.microsoft.com/office/powerpoint/2010/main"
    for sect_lst in prs_elem.iter(QName(ns14, "sectionLst")):
        sect_lst.getparent().remove(sect_lst)
